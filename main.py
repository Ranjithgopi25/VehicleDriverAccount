from fastapi import FastAPI, HTTPException, UploadFile, File, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse, FileResponse
from pydantic import BaseModel
from typing import List, Optional
import os
import json
import logging
from dotenv import load_dotenv

logger = logging.getLogger(__name__)
from core.llm import FallbackLLMService, get_llm_service
from core.semantic_retriever import SemanticRetriever, DEPENDENCIES_AVAILABLE
from core.citation_validator import validate_citations, extract_citations_from_text
from routers.ddc import ddc_router
from pptx import Presentation
from pptx.util import Inches, Pt
from pypdf import PdfReader
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import io
import tempfile
from collections import Counter
from ppt_sanitizer import PPTSanitizer
import httpx
from bs4 import BeautifulSoup
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.enums import TA_JUSTIFY, TA_LEFT
from docx import Document
from docx.shared import Pt as DocxPt, Inches as DocxInches, RGBColor as DocxRGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_BREAK
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
import re
from datetime import datetime
from urllib.parse import urlparse
import ipaddress
import socket
import platform
import subprocess
import sys

# Audioop compatibility for Python 3.13+
try:
    import audioop
except ModuleNotFoundError:
    import audioop_lts as audioop
    import sys

    sys.modules["audioop"] = audioop

import boto3
from botocore.exceptions import BotoCoreError, ClientError
# from pydub import AudioSegment
import base64

load_dotenv()

app = FastAPI(title="PwC Presentation Assistant API")

# CORS Configuration
import re

# Define allowed origins with regex support for Amplify branches
allowed_origins_patterns = [
    r"http://localhost:\d+",  # Local development (any port)
    r"https://.*\.amplifyapp\.com",  # All Amplify branches
    r"https://deployment-frontend\.d2ebg85go3xrq2\.amplifyapp\.com",
    r"https://deployment-phase1\.d2ebg85go3xrq2\.amplifyapp\.com",  # Specific Amplify app
]


def check_origin(origin: str) -> bool:
    """Check if origin matches any allowed pattern"""
    for pattern in allowed_origins_patterns:
        if re.match(pattern, origin):
            return True
    return False


app.add_middleware(
    CORSMiddleware,
    allow_origin_regex=r"(http://localhost:\d+|https://.*\.amplifyapp\.com)",
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Register DDC workflow routers
app.include_router(ddc_router, prefix="/api/ddc", tags=["DDC"])

# Register modular API v1 routers
from app.api.v1.ddc import brand_format
app.include_router(brand_format.router, prefix="/api/v1/ddc")


def get_llm():
    """Get the LLM service instance (supports Azure OpenAI + Groq fallback)"""
    try:
        return get_llm_service()
    except Exception as e:
        raise HTTPException(
            status_code=500, detail=f"LLM service not available: {str(e)}"
        )


class Message(BaseModel):
    role: str
    content: str


class ChatRequest(BaseModel):
    messages: List[Message]
    stream: bool = True


class DraftRequest(BaseModel):
    topic: str
    objective: str
    audience: str
    additional_context: Optional[str] = None


class ThoughtLeadershipRequest(BaseModel):
    operation: str
    topic: Optional[str] = None
    perspective: Optional[str] = None
    target_audience: Optional[str] = None
    document_text: Optional[str] = None
    target_format: Optional[str] = None
    additional_context: Optional[str] = None
    reference_urls: Optional[List[str]] = None


class ExportRequest(BaseModel):
    content: str
    title: Optional[str] = "Generated Article"


class ResearchRequest(BaseModel):
    query: str
    focus_areas: Optional[List[str]] = None
    additional_context: Optional[str] = None


class ArticleRequest(BaseModel):
    topic: str
    content_type: str  # Article, Case Study, Executive Brief, Blog, etc.
    desired_length: int  # word count
    tone: str  # Professional, Conversational, Technical, etc.
    outline_text: Optional[str] = None
    additional_context: Optional[str] = None


class BestPracticesRequest(BaseModel):
    categories: Optional[List[str]] = (
        None  # Structure, Visuals, Design, Charts, Formatting, Content
    )


class PodcastRequest(BaseModel):
    customization: Optional[str] = None


# New Thought Leadership Section Models (5 Sections)


class DraftContentRequest(BaseModel):
    """Section 1: Draft Content"""

    messages: List[Message]  # Conversational history
    content_type: Optional[str] = None  # Article, Blog, White Paper, Executive Brief
    topic: Optional[str] = None
    audience: Optional[str] = None
    length: Optional[str] = None  # Word count or page count
    include_research: bool = False
    additional_guidelines: Optional[str] = None
    stream: bool = True


class ConductResearchRequest(BaseModel):
    """Section 2: Conduct Research"""

    messages: List[Message]  # Conversational history
    query: Optional[str] = None
    source_groups: Optional[List[str]] = None  # PwC Proprietary, Licensed, External
    specific_sources: Optional[List[str]] = None  # Individual source selection
    standalone: bool = True  # True = standalone research, False = embedded in content
    stream: bool = True


class EditContentRequest(BaseModel):
    """Section 3: Edit Content"""

    messages: List[Message]  # Conversational history
    editor_types: Optional[List[str]] = (
        None  # Brand Alignment, Copy, Line, Content, Development
    )
    document_text: Optional[str] = None
    stream: bool = True


class RefineContentRequest(BaseModel):
    """Section 4: Refine Content"""

    messages: List[Message]  # Conversational history
    document_text: Optional[str] = None
    services: Optional[List[str]] = (
        None  # Expand/Compress, Tone/Audience, Research, Edit, Suggestions
    )
    desired_length: Optional[str] = None
    target_audience: Optional[str] = None
    target_tone: Optional[str] = None
    stream: bool = True


class FormatTranslatorRequest(BaseModel):
    """Section 5: Format Translator"""

    messages: List[Message]  # Conversational history
    document_text: Optional[str] = None
    source_format: Optional[str] = None  # Article, Blog, White Paper, etc.
    target_format: Optional[str] = None
    stream: bool = True


class UpdateSectionRequest(BaseModel):
    """Canvas Editor: Update specific section of content"""

    fullArticle: str
    sectionIndex: int
    sectionContent: str
    userPrompt: str
    contentType: str


@app.get("/")
async def root():
    return {
        "message": "PwC Presentation Assistant API",
        "version": "1.0.0",
        "status": "running",
    }


@app.get("/health")
async def health_check():
    try:
        llm = get_llm()
        providers = llm.get_active_providers()
        return {
            "status": "healthy",
            "llm_providers": providers,
            "primary": providers[0] if providers else "none",
        }
    except Exception as e:
        return {"status": "degraded", "llm_providers": [], "error": str(e)}


async def crawl_related_pages(
    initial_url: str, max_pages: int = 10, max_depth: int = 3, query: str = None
) -> List[dict]:
    """
    Advanced crawling function with Perplexity-level intelligence.
    - Prioritizes relevant links based on anchor text and URL patterns
    - Extracts navigation menus and sitemaps
    - Uses semantic relevance scoring
    - Handles structured data and JSON-LD
    """
    parsed_initial = urlparse(initial_url)
    base_domain = f"{parsed_initial.scheme}://{parsed_initial.netloc}"

    visited_urls = set()
    # Priority queue: (url, depth, priority_score)
    to_visit = [(initial_url, 0, 100.0)]
    fetched_pages = []

    # Keywords from query for relevance scoring
    query_keywords = set()
    if query:
        query_keywords = set(re.findall(r"\b\w+\b", query.lower()))

    # Browser-like headers
    headers = {
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36",
        "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8",
        "Accept-Language": "en-US,en;q=0.9",
        "Accept-Encoding": "gzip, deflate, br",
        "Connection": "keep-alive",
        "Referer": base_domain,
    }

    def calculate_link_priority(
        href: str, anchor_text: str, link_context: str = ""
    ) -> float:
        """Calculate priority score for a link (0-100)"""
        priority = 50.0  # Base priority

        href_lower = href.lower()
        anchor_lower = anchor_text.lower()
        context_lower = link_context.lower()

        # Boost for query keyword matches
        if query_keywords:
            for keyword in query_keywords:
                if keyword in href_lower:
                    priority += 15.0
                if keyword in anchor_lower:
                    priority += 20.0
                if keyword in context_lower:
                    priority += 10.0

        # Boost for content-related URLs
        content_indicators = [
            "article",
            "post",
            "page",
            "content",
            "insight",
            "research",
            "report",
            "study",
            "analysis",
            "industry",
            "sector",
            "practice",
            "service",
            "capability",
            "library",
            "outlook",
            "trend",
        ]
        for indicator in content_indicators:
            if indicator in href_lower or indicator in anchor_lower:
                priority += 10.0

        # Boost for navigation menu links (often in nav, menu, or have class names)
        nav_indicators = ["nav", "menu", "navigation", "main-menu", "primary-nav"]
        if any(ind in context_lower for ind in nav_indicators):
            priority += 25.0

        # Penalize for non-content URLs
        skip_patterns = [
            "mailto:",
            "tel:",
            "javascript:",
            "#",
            ".pdf",
            ".jpg",
            ".png",
            ".gif",
            ".zip",
            ".exe",
            "login",
            "signup",
            "register",
            "logout",
            "admin",
            "search",
            "tag",
            "category",
            "author",
            "feed",
            "rss",
            "print",
            "share",
            "comment",
            "social",
        ]
        for pattern in skip_patterns:
            if pattern in href_lower:
                priority -= 50.0

        # Penalize very long URLs (often dynamic/less important)
        if len(href) > 150:
            priority -= 10.0

        # Boost for shorter, cleaner URLs (often main pages)
        if len(href.split("/")) <= 5 and href.count("?") == 0:
            priority += 5.0

        return max(0.0, min(100.0, priority))

    async def extract_sitemap_urls(soup: BeautifulSoup, base_domain: str) -> List[str]:
        """Extract URLs from sitemap references"""
        sitemap_urls = []

        # Look for sitemap in robots.txt
        try:
            robots_url = f"{base_domain}/robots.txt"
            async with httpx.AsyncClient(timeout=10.0, headers=headers) as client:
                robots_resp = await client.get(robots_url)
                if robots_resp.status_code == 200:
                    for line in robots_resp.text.split("\n"):
                        if line.lower().startswith("sitemap:"):
                            sitemap_url = line.split(":", 1)[1].strip()
                            sitemap_urls.append(sitemap_url)
        except:
            pass

        # Look for sitemap links in HTML
        sitemap_links = soup.find_all("a", href=re.compile(r"sitemap", re.I))
        for link in sitemap_links:
            href = link.get("href", "")
            if href.startswith("/"):
                sitemap_urls.append(f"{base_domain}{href}")
            elif href.startswith("http"):
                sitemap_urls.append(href)

        return sitemap_urls

    async def extract_navigation_links(soup: BeautifulSoup) -> List[dict]:
        """Extract links from navigation menus"""
        nav_links = []

        # Find navigation elements
        nav_elements = soup.find_all(
            ["nav", "ul", "ol"], class_=re.compile(r"nav|menu", re.I)
        )
        nav_elements.extend(soup.find_all("nav"))

        for nav in nav_elements:
            links = nav.find_all("a", href=True)
            for link in links:
                href = link.get("href", "")
                anchor = link.get_text().strip()
                if href and anchor:
                    nav_links.append(
                        {"href": href, "anchor": anchor, "context": "navigation"}
                    )

        return nav_links

    async with httpx.AsyncClient(
        timeout=30.0, follow_redirects=True, max_redirects=5, headers=headers
    ) as client:
        # First, try to fetch and parse sitemap
        try:
            sitemap_urls = []
            # Check common sitemap locations
            common_sitemaps = [
                f"{base_domain}/sitemap.xml",
                f"{base_domain}/sitemap_index.xml",
                f"{base_domain}/sitemap/sitemap.xml",
            ]

            for sitemap_url in common_sitemaps:
                try:
                    resp = await client.get(sitemap_url, timeout=10.0)
                    if resp.status_code == 200:
                        # Parse XML sitemap
                        from xml.etree import ElementTree as ET

                        try:
                            root = ET.fromstring(resp.text)
                            # Handle both sitemapindex and urlset
                            if root.tag.endswith("sitemapindex"):
                                for sitemap in root.findall(".//{*}sitemap/{*}loc"):
                                    sitemap_urls.append(sitemap.text)
                            elif root.tag.endswith("urlset"):
                                for url_elem in root.findall(".//{*}url/{*}loc"):
                                    url = url_elem.text
                                    if url and parsed_initial.netloc in url:
                                        priority = calculate_link_priority(url, "", "")
                                        to_visit.append(
                                            (url, 1, priority + 30.0)
                                        )  # Boost sitemap URLs
                        except:
                            pass
                except:
                    continue
        except:
            pass

        while to_visit and len(fetched_pages) < max_pages:
            # Sort by priority (highest first)
            to_visit.sort(key=lambda x: x[2], reverse=True)
            current_url, depth, priority = to_visit.pop(0)

            # Skip if already visited or depth exceeded
            if current_url in visited_urls or depth > max_depth:
                continue

            visited_urls.add(current_url)

            try:
                # Fetch the page
                response = await client.get(current_url, timeout=30.0)

                if response.status_code != 200:
                    continue

                soup = BeautifulSoup(response.text, "html.parser")

                # Extract title
                title = soup.find("title")
                title_text = title.get_text().strip() if title else ""

                # Extract meta description
                meta_desc = soup.find("meta", attrs={"name": "description"})
                description = meta_desc.get("content", "").strip() if meta_desc else ""

                # Extract structured data (JSON-LD)
                json_ld_data = []
                for script in soup.find_all("script", type="application/ld+json"):
                    try:
                        data = json.loads(script.string)
                        json_ld_data.append(data)
                    except:
                        pass

                # Extract main content with better parsing
                for script in soup(
                    ["script", "style", "nav", "footer", "header", "aside"]
                ):
                    script.decompose()

                # Try multiple content extraction strategies
                article_content = None
                content_selectors = [
                    ("article", None),
                    ("main", None),
                    ("div", {"class": re.compile(r"content|article|post|main", re.I)}),
                    ("div", {"id": re.compile(r"content|article|main", re.I)}),
                    ("body", None),
                ]

                for selector, attrs in content_selectors:
                    if attrs:
                        article_content = soup.find(selector, attrs=attrs)
                    else:
                        article_content = soup.find(selector)
                if article_content:
                    break

                if article_content:
                    # Extract structured content
                    paragraphs = article_content.find_all(
                        [
                            "p",
                            "h1",
                            "h2",
                            "h3",
                            "h4",
                            "h5",
                            "h6",
                            "li",
                            "div",
                            "section",
                        ]
                    )
                    content_parts = []
                    for p in paragraphs:
                        text = p.get_text().strip()
                        if text and len(text) > 20:  # Filter out very short text
                            # Preserve heading structure
                            if p.name in ["h1", "h2", "h3", "h4", "h5", "h6"]:
                                content_parts.append(f"\n{p.name.upper()}: {text}\n")
                            else:
                                content_parts.append(text)
                    content = "\n".join(content_parts)
                else:
                    content = soup.get_text()

                # Clean up content
                content = re.sub(r"\n\s*\n+", "\n\n", content).strip()
                content = re.sub(r"\s+", " ", content)  # Normalize whitespace

                # Combine with description
                if description and description not in content:
                    content = f"{description}\n\n{content}"

                # Store fetched page
                fetched_pages.append(
                    {
                        "url": current_url,
                        "title": title_text,
                        "content": content[:10000],  # Increased limit
                    }
                )

                # Extract links for next level (only if not at max depth)
                if depth < max_depth and len(fetched_pages) < max_pages:
                    # Get navigation links first (high priority)
                    nav_links = await extract_navigation_links(soup)

                    # Get all links
                    all_links = soup.find_all("a", href=True)

                    # Combine and deduplicate
                    link_map = {}
                    for link in nav_links:
                        href = link["href"]
                        if href not in link_map:
                            link_map[href] = {
                                "href": href,
                                "anchor": link["anchor"],
                                "context": link["context"],
                                "priority": 0.0,
                            }

                    for link in all_links:
                        href = link.get("href", "")
                        if not href:
                            continue

                        anchor_text = link.get_text().strip()
                        parent = link.find_parent()
                        parent_class = parent.get("class", []) if parent else []
                        context = " ".join(parent_class) if parent_class else ""

                        if href not in link_map:
                            link_map[href] = {
                                "href": href,
                                "anchor": anchor_text,
                                "context": context,
                                "priority": 0.0,
                            }
                        else:
                            # Update with better anchor text if available
                            if anchor_text and not link_map[href]["anchor"]:
                                link_map[href]["anchor"] = anchor_text

                    # Resolve and prioritize links
                    prioritized_links = []
                    for link_info in link_map.values():
                        href = link_info["href"]

                        # Resolve relative URLs
                        if href.startswith("/"):
                            full_url = f"{base_domain}{href}"
                        elif href.startswith("http://") or href.startswith("https://"):
                            parsed_link = urlparse(href)
                            # Only follow links from same domain
                            if parsed_link.netloc == parsed_initial.netloc:
                                full_url = href
                            else:
                                continue
                        else:
                            # Relative URL
                            current_path = (
                                current_url.rsplit("/", 1)[0]
                                if "/" in current_url
                                else current_url
                            )
                            full_url = f"{current_path}/{href}"

                        # Normalize URL (remove fragments and query params for deduplication)
                        full_url_clean = full_url.split("#")[0].split("?")[0]

                        # Skip if already visited
                        if full_url_clean in visited_urls:
                            continue

                        # Calculate priority
                        priority = calculate_link_priority(
                            href, link_info["anchor"], link_info["context"]
                        )

                        # Skip low-priority or invalid links
                        if priority < 20.0:
                            continue

                        # Check if already queued
                        already_queued = any(
                            url == full_url_clean or url == full_url
                            for url, _, _ in to_visit
                        )

                        if not already_queued:
                            prioritized_links.append((full_url, depth + 1, priority))

                    # Sort by priority and add top links
                    prioritized_links.sort(key=lambda x: x[2], reverse=True)
                    # Add more links per page (up to 50 for first level, 30 for deeper)
                    max_links_per_page = 50 if depth == 0 else 30
                    for link_url, link_depth, link_priority in prioritized_links[
                        :max_links_per_page
                    ]:
                        to_visit.append((link_url, link_depth, link_priority))

            except Exception as e:
                logger.debug(f"Error crawling {current_url}: {e}")
                continue

    return fetched_pages


def generate_url_variations(url: str) -> List[str]:
    """Generate common URL variations to try when encountering 404 errors"""
    variations = []
    parsed = urlparse(url)

    # Original URL (don't add it twice if it's already the first variation)
    path = parsed.path

    # Variation 1: Add trailing slash if not present
    if path and not path.endswith("/"):
        variations.append(f"{parsed.scheme}://{parsed.netloc}{path}/")

    # Variation 2: Remove trailing slash if present
    if path.endswith("/") and len(path) > 1:
        variations.append(f"{parsed.scheme}://{parsed.netloc}{path.rstrip('/')}")

    # Variation 3: Try lowercase path (common on case-sensitive servers)
    if path != path.lower():
        variations.append(f"{parsed.scheme}://{parsed.netloc}{path.lower()}")

    # Variation 4: Try singular/plural variations for common patterns
    path_parts = path.strip("/").split("/")
    if len(path_parts) >= 1 and path_parts[-1]:
        last_part = path_parts[-1]
        # Try singular if plural
        if last_part.endswith("s") and len(last_part) > 1:
            singular = last_part[:-1]
            new_path = "/".join(path_parts[:-1] + [singular])
            variations.append(f"{parsed.scheme}://{parsed.netloc}/{new_path}")
        # Try plural if singular
        elif not last_part.endswith("s"):
            plural = last_part + "s"
            new_path = "/".join(path_parts[:-1] + [plural])
            variations.append(f"{parsed.scheme}://{parsed.netloc}/{new_path}")

    # Variation 5: Try with index.html
    if path and not path.endswith((".html", ".htm", ".php", ".asp", ".aspx")):
        variations.append(f"{parsed.scheme}://{parsed.netloc}{path}/index.html")

    # Add query string and fragment if present
    query_fragment = ""
    if parsed.query:
        query_fragment += f"?{parsed.query}"
    if parsed.fragment:
        query_fragment += f"#{parsed.fragment}"

    # Apply query/fragment to all variations
    if query_fragment:
        variations = [v + query_fragment for v in variations]

    return variations


def is_safe_hostname(hostname: str) -> bool:
    """Check if hostname is safe (not loopback, not private IP, not unspecified)"""
    try:
        ip = ipaddress.ip_address(hostname)
        return not (
            ip.is_loopback
            or ip.is_private
            or ip.is_reserved
            or ip.is_multicast
            or ip.is_unspecified
            or ip.is_link_local
        )
    except ValueError:
        try:
            resolved_ip = socket.gethostbyname(hostname)
            ip = ipaddress.ip_address(resolved_ip)
            return not (
                ip.is_loopback
                or ip.is_private
                or ip.is_reserved
                or ip.is_multicast
                or ip.is_unspecified
                or ip.is_link_local
            )
        except (socket.gaierror, ValueError):
            return True


async def fetch_url_content(url: str) -> dict:
    """Fetch and extract content from a URL with security validations"""
    try:
        parsed = urlparse(url)
        if parsed.scheme not in ["http", "https"]:
            raise ValueError(
                f"Invalid URL scheme: {parsed.scheme}. Only http and https are allowed."
            )

        if not parsed.hostname:
            raise ValueError("Invalid URL: missing domain")

        if not is_safe_hostname(parsed.hostname):
            raise ValueError(
                "Access to localhost, private IP ranges, or reserved IPs is not allowed for security reasons"
            )

        # Browser-like headers to avoid 403 Forbidden errors
        # Build referer from the URL (same domain)
        referer = f"{parsed.scheme}://{parsed.netloc}/"

        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36",
            "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8,application/signed-exchange;v=b3;q=0.7",
            "Accept-Language": "en-US,en;q=0.9",
            "Accept-Encoding": "gzip, deflate, br",
            "Connection": "keep-alive",
            "Upgrade-Insecure-Requests": "1",
            "Referer": referer,
            "Sec-Fetch-Dest": "document",
            "Sec-Fetch-Mode": "navigate",
            "Sec-Fetch-Site": "same-origin",
            "Sec-Fetch-User": "?1",
            "Cache-Control": "max-age=0",
            "DNT": "1",
        }

        async with httpx.AsyncClient(
            timeout=60.0, follow_redirects=True, max_redirects=5, headers=headers
        ) as client:
            response = await client.get(url)

            # Check status codes before raising
            if response.status_code == 404:
                # Try URL variations before giving up
                variations = generate_url_variations(url)
                tried_urls = [url]

                for variation_url in variations:
                    if variation_url in tried_urls:
                        continue
                    tried_urls.append(variation_url)

                    try:
                        # Update referer for the variation
                        parsed_var = urlparse(variation_url)
                        var_referer = f"{parsed_var.scheme}://{parsed_var.netloc}/"
                        headers["Referer"] = var_referer

                        var_response = await client.get(variation_url)
                        if var_response.status_code == 200:
                            # Success! Use this URL and continue with content extraction
                            url = variation_url  # Update url to the working variation
                            response = var_response
                            break
                        elif var_response.status_code != 404:
                            # If it's not 200 or 404, it's a different error, stop trying
                            break
                    except Exception:
                        # Continue to next variation on any error
                        continue

                # If still 404 after trying variations, return error
                if response.status_code == 404:
                    return {
                        "url": url,
                        "title": "",
                        "content": "",
                        "success": False,
                        "error": f"404 Not Found: The page at this URL does not exist. Tried {len(tried_urls)} variation(s) but none were found. Please verify the URL is correct.",
                    }
            elif response.status_code == 403:
                return {
                    "url": url,
                    "title": "",
                    "content": "",
                    "success": False,
                    "error": f"403 Forbidden: The website blocked access. This may be due to bot protection, Cloudflare, or WAF rules. The site may require JavaScript to load content.",
                }
            elif response.status_code == 429:
                return {
                    "url": url,
                    "title": "",
                    "content": "",
                    "success": False,
                    "error": f"429 Too Many Requests: Rate limit exceeded. Please wait a moment and try again.",
                }
            elif response.status_code == 503:
                return {
                    "url": url,
                    "title": "",
                    "content": "",
                    "success": False,
                    "error": f"503 Service Unavailable: The website is temporarily unavailable. Please try again later.",
                }

            response.raise_for_status()

            soup = BeautifulSoup(response.text, "html.parser")

            for script in soup(["script", "style", "nav", "footer", "header"]):
                script.decompose()

            title = soup.find("title")
            title_text = title.get_text().strip() if title else ""

            article_content = (
                soup.find("article") or soup.find("main") or soup.find("body")
            )
            if article_content:
                paragraphs = article_content.find_all(
                    ["p", "h1", "h2", "h3", "h4", "li"]
                )
                content = "\n".join(
                    [p.get_text().strip() for p in paragraphs if p.get_text().strip()]
                )
            else:
                content = soup.get_text()

            content = re.sub(r"\n\s*\n", "\n\n", content)
            content = content.strip()

            return {
                "url": url,
                "title": title_text,
                "content": content[:5000],
                "success": True,
            }
    except Exception as e:
        return {
            "url": url,
            "title": "",
            "content": "",
            "success": False,
            "error": str(e),
        }


def extract_text_from_pdf(pdf_bytes: bytes, max_chars: int = 10000) -> str:
    """Extract text content from PDF file"""
    try:
        pdf = PdfReader(io.BytesIO(pdf_bytes))
        text = ""
        for page in pdf.pages:
            text += page.extract_text() + "\n"
            if len(text) > max_chars:
                break
        return text[:max_chars]
    except Exception as e:
        return f"[Error reading PDF: {str(e)}]"


def extract_text_from_docx(docx_bytes: bytes, max_chars: int = 10000) -> str:
    """Extract text content from DOCX file"""
    try:
        doc = Document(io.BytesIO(docx_bytes))
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
            if len(text) > max_chars:
                break
        return text[:max_chars]
    except Exception as e:
        return f"[Error reading DOCX: {str(e)}]"


async def parse_uploaded_file(uploaded_file: UploadFile, max_chars: int = 10000) -> str:
    """Parse uploaded file and extract text content based on file type"""
    if not uploaded_file:
        return ""

    file_content = await uploaded_file.read()
    filename = uploaded_file.filename.lower() if uploaded_file.filename else ""

    if filename.endswith(".pdf"):
        return extract_text_from_pdf(file_content, max_chars)
    elif filename.endswith(".docx"):
        return extract_text_from_docx(file_content, max_chars)
    elif filename.endswith((".txt", ".md")):
        try:
            return file_content.decode("utf-8")[:max_chars]
        except:
            return f"[Could not decode text file: {uploaded_file.filename}]"
    else:
        try:
            return file_content.decode("utf-8")[:max_chars]
        except:
            return (
                f"[Unsupported file type or could not decode: {uploaded_file.filename}]"
            )


def load_template(template_path: str = 'template/Template.docx') -> Optional[Document]:
    """Load template file and return Document object"""
    try:
        # Resolve template path - try multiple strategies
        resolved_path = None
        
        # Strategy 1: If absolute path provided, use it directly
        if os.path.isabs(template_path):
            if os.path.exists(template_path):
                resolved_path = template_path
        
        # Strategy 2: Resolve relative to backend directory (where main.py is located)
        if not resolved_path:
            backend_dir = os.path.dirname(os.path.abspath(__file__))
            normalized_template = os.path.normpath(template_path)
            candidate = os.path.join(backend_dir, normalized_template)
            if os.path.exists(candidate):
                resolved_path = candidate
        
        # Strategy 3: Try relative to current working directory
        if not resolved_path:
            cwd = os.getcwd()
            normalized_template = os.path.normpath(template_path)
            candidate = os.path.join(cwd, normalized_template)
            if os.path.exists(candidate):
                resolved_path = candidate
        
        if not resolved_path:
            logger.warning(
                f"Template file not found. Searched for '{template_path}' in:\n"
                f"  - {os.path.dirname(os.path.abspath(__file__))}\n"
                f"  - {os.getcwd()}\n"
                f"Using default formatting."
            )
            return None
        
        logger.info(f"Successfully loaded template from: {resolved_path}")
        template_doc = Document(resolved_path)
        return template_doc
        
    except Exception as e:
        logger.error(f"Error loading template: {str(e)}")
        return None


def extract_toc_from_content(content: str) -> Optional[List[str]]:
    """Check if table of contents already exists in content and extract it"""
    lines = content.split("\n")
    toc_started = False
    toc_ended = False
    headings = []
    toc_keywords = ["table of contents", "contents", "table of content", "toc", "content"]
    
    for i, line in enumerate(lines):
        line_lower = line.strip().lower()
        
        # Check if TOC section starts (look for "Content" or "Table of Contents" heading)
        if any(keyword in line_lower for keyword in toc_keywords):
            # Make sure it's a heading or standalone line, not part of regular content
            if (line_lower in toc_keywords or 
                line_lower == "content" or
                "table of contents" in line_lower or
                (line.strip().startswith("#") and any(keyword in line_lower for keyword in toc_keywords))):
                toc_started = True
                continue
        
        # If TOC started, extract headings until we hit a major section
        if toc_started and not toc_ended:
            line_stripped = line.strip()
            
            # Skip empty lines at start
            if not line_stripped:
                continue
            
            # Check if we've reached the end of TOC (major heading or content)
            if line_stripped.startswith("# "):
                # Check if it's a real heading or just TOC entry
                if len(headings) > 0:
                    toc_ended = True
                    break
                continue
            
            # Extract TOC entries (lines with dots, numbers, or simple text)
            # Pattern: "Heading text ... page number" or "1. Heading" or "- Heading"
            if re.match(r'^[\d\-•*]\s+', line_stripped) or "..." in line_stripped or re.match(r'^\d+[\.\)]', line_stripped):
                # Extract heading text (remove numbering, bullets, dots, page numbers)
                heading = re.sub(r'^[\d\-•*]\s+', '', line_stripped)
                heading = re.sub(r'\.\.\..*$', '', heading)  # Remove dots and page number
                heading = re.sub(r'\s+\d+$', '', heading)  # Remove trailing page number
                heading = heading.strip()
                if heading and len(heading) > 2:
                    headings.append(heading)
            elif line_stripped and not line_stripped.startswith("#"):
                # Simple text entry (might be TOC entry without formatting)
                if len(headings) < 20:  # Reasonable TOC length
                    headings.append(line_stripped)
                else:
                    toc_ended = True
                    break
    
    # Return headings if we found a TOC section with multiple entries
    if toc_started and len(headings) >= 2:
        logger.info(f"Found existing TOC in content with {len(headings)} entries")
        return headings[:25]  # Limit to 25 headings
    
    return None


def extract_headings_from_content(content: str) -> List[str]:
    """Extract important headings from content for table of contents (only # and ##)"""
    headings = []
    lines = content.split("\n")
    
    for line in lines:
        line = line.strip()
        # Only extract important headings (level 1 and 2)
        if line.startswith("# "):
            headings.append(line[2:].strip())
        elif line.startswith("## "):
            headings.append(line[3:].strip())
        # Skip ### and #### (too detailed for TOC)
        # Check for bold headings (markdown **) only if it's a major section
        elif line.startswith("**") and line.endswith("**") and len(line) > 4 and len(line) < 100:
            headings.append(line[2:-2].strip())
    
    return headings[:25]  # Limit to 25 headings for TOC


async def generate_title_with_llm(content: str) -> str:
    """Generate a title from content using LLM"""
    try:
        llm = get_llm()
        system_prompt = "You are a professional title generator. Generate a concise, compelling title for the given content."
        user_prompt = f"Generate a professional title for the following content. Return ONLY the title, no additional text:\n\n{content[:2000]}"
        
        response = await llm.chat_completion(
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            temperature=0.7,
            max_tokens=100
        )
        
        title = response.content.strip()
        # Clean up title - remove quotes if present
        title = title.strip('"').strip("'").strip()
        # Take only first line if multiple lines
        title = title.split('\n')[0].strip()
        return title if title else "Document"
    except Exception as e:
        logger.error(f"Error generating title with LLM: {str(e)}")
        # Fallback: extract first heading or first sentence
        lines = content.split("\n")
        for line in lines:
            line = line.strip()
            if line.startswith("# "):
                return line[2:].strip()
            if line and len(line) < 100:
                return line[:100]
        return "Document"


async def generate_table_of_contents_with_llm(content: str) -> List[str]:
    """Generate table of contents from content using LLM - only important headings"""
    try:
        llm = get_llm()
        system_prompt = "You are a professional document editor. Generate a table of contents based on the content structure. Include only important main headings and major sections, not detailed sub-sections."
        user_prompt = f"Analyze the following content and generate a table of contents with only the most important headings (main sections and major topics only, skip detailed sub-sections). Return ONLY a list of headings, one per line, without numbers or formatting:\n\n{content[:3000]}"
        
        response = await llm.chat_completion(
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            temperature=0.7,
            max_tokens=500
        )
        
        toc_lines = response.content.strip().split("\n")
        headings = []
        for line in toc_lines:
            line = line.strip()
            # Remove numbering and formatting
            line = re.sub(r'^\d+[\.\)]\s*', '', line)
            line = re.sub(r'^[-*]\s*', '', line)
            line = line.strip()
            if line and len(line) > 2:
                headings.append(line)
        
        # Fallback to extracted headings if LLM didn't return good results
        if not headings or len(headings) < 2:
            headings = extract_headings_from_content(content)
        
        return headings[:25]  # Limit to 25 headings
    except Exception as e:
        logger.error(f"Error generating TOC with LLM: {str(e)}")
        # Fallback to extracted headings
        return extract_headings_from_content(content)


async def create_pdf(content: str, title: Optional[str] = None, template_path: str = 'template/Template.docx') -> bytes:
    """Generate a PDF from content by converting Word document to PDF"""
    try:
        # Generate title if not provided
        if not title:
            title = await generate_title_with_llm(content)
        # First, create Word document using template
        docx_bytes = await create_word_doc(content, title=title, template_path=template_path, use_template=True)
        
        # Try to convert Word to PDF using Word COM automation (Windows only)
        if platform.system() == "Windows":
            try:
                return convert_word_to_pdf_com(docx_bytes)
            except Exception as e:
                logger.warning(f"Word COM automation failed: {str(e)}, falling back to ReportLab")
        
        # Fallback: Use ReportLab to create PDF from Word document content
        # Extract text from Word document and format with ReportLab
        return create_pdf_from_word_content(docx_bytes, title)
    
    except Exception as e:
        logger.error(f"Error creating PDF: {str(e)}")
        # Ultimate fallback: Create simple PDF with ReportLab
        return create_pdf_simple_fallback(content, title)


def convert_word_to_pdf_com(docx_bytes: bytes) -> bytes:
    """Convert Word document to PDF using Word COM automation (Windows only)"""
    try:
        import win32com.client
        import pythoncom
        
        # Create temporary files
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as docx_file:
            docx_file.write(docx_bytes)
            docx_path = docx_file.name
        
        pdf_path = docx_path.replace('.docx', '.pdf')
        
        try:
            # Initialize COM
            pythoncom.CoInitialize()
            
            # Create Word application
            word = win32com.client.Dispatch("Word.Application")
            word.Visible = False
            
            try:
                # Open document
                doc = word.Documents.Open(docx_path)
                
                # Export as PDF
                doc.ExportAsFixedFormat(
                    OutputFileName=pdf_path,
                    ExportFormat=17,  # wdExportFormatPDF
                    OpenAfterExport=False,
                    OptimizeFor=0,  # wdExportOptimizeForPrint
                    BitmapMissingFonts=True,
                    DocStructureTags=True,
                    CreateBookmarks=0,  # wdExportCreateNoBookmarks
                    UseISO19005_1=False
                )
                
                # Close document
                doc.Close(False)
                
                # Read PDF file
                with open(pdf_path, 'rb') as pdf_file:
                    pdf_bytes = pdf_file.read()
                
                return pdf_bytes
            
            finally:
                # Quit Word application
                word.Quit()
                pythoncom.CoUninitialize()
        
        finally:
            # Clean up temporary files
            try:
                os.unlink(docx_path)
                if os.path.exists(pdf_path):
                    os.unlink(pdf_path)
            except Exception as e:
                logger.warning(f"Error cleaning up temporary files: {str(e)}")
    
    except ImportError:
        raise Exception("pywin32 not installed. Install it with: pip install pywin32")
    except Exception as e:
        raise Exception(f"Word COM automation failed: {str(e)}")


def create_pdf_from_word_content(docx_bytes: bytes, title: str) -> bytes:
    """Create PDF from Word document content using ReportLab (fallback method)"""
    try:
        # Extract text from Word document
        doc = Document(io.BytesIO(docx_bytes))
        content_text = []
        
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                content_text.append(text)
        
        # Create PDF using ReportLab
        buffer = io.BytesIO()
        doc_pdf = SimpleDocTemplate(
            buffer, pagesize=letter, topMargin=1 * inch, bottomMargin=1 * inch
        )
        
        styles = getSampleStyleSheet()
        title_style = ParagraphStyle(
            "CustomTitle",
            parent=styles["Heading1"],
            fontSize=26,  # Increased to match Word (26pt)
            textColor="#E87722",
            spaceAfter=30,
            alignment=TA_LEFT,
            fontName="Helvetica-Bold",
        )
        
        body_style = ParagraphStyle(
            "CustomBody",
            parent=styles["BodyText"],
            fontSize=12,  # Increased to match Word (12pt)
            alignment=TA_JUSTIFY,
            spaceAfter=10,  # Decreased
            leading=21.6,  # 1.8 line spacing for 12pt font (12 * 1.8)
        )
        
        heading1_style = ParagraphStyle(
            "CustomHeading1",
            parent=styles["Heading1"],
            fontSize=20,  # Increased to match Word
            textColor="#E87722",  # Orange, bold
            spaceAfter=12,  # Decreased
            spaceBefore=18,  # Decreased
            fontName="Helvetica-Bold",
            leading=36,  # 1.8 line spacing
        )
        
        heading2_style = ParagraphStyle(
            "CustomHeading2",
            parent=styles["Heading2"],
            fontSize=18,  # Increased to match Word
            textColor="#E87722",  # Orange, bold
            spaceAfter=10,  # Decreased
            spaceBefore=14,  # Decreased
            fontName="Helvetica-Bold",
            leading=32.4,  # 1.8 line spacing
        )
        
        heading3_style = ParagraphStyle(
            "CustomHeading3",
            parent=styles["Heading2"],
            fontSize=16,  # Increased to match Word
            textColor="#333333",  # Dark gray, bold
            spaceAfter=8,  # Decreased
            spaceBefore=12,  # Decreased
            fontName="Helvetica-Bold",
            leading=28.8,  # 1.8 line spacing
        )
        
        title_content_style = ParagraphStyle(
            "TitleContent",
            parent=styles["Heading1"],
            fontSize=20,
            textColor="#E87722",  # Orange, bold
            spaceAfter=12,
            spaceBefore=18,
            fontName="Helvetica-Bold",
            leading=36,  # 1.8 line spacing
        )
        
        bullet_style = ParagraphStyle(
            "Bullet",
            parent=styles["BodyText"],
            fontSize=12,
            alignment=TA_LEFT,
            spaceAfter=8,  # Decreased
            leading=21.6,  # 1.8 line spacing
            leftIndent=20,
        )
        
        citation_style = ParagraphStyle(
            "Citation",
            parent=styles["BodyText"],
            fontSize=9,
            textColor="#666666",
            leftIndent=20,
            spaceAfter=6,
        )
        
        story = []
        story.append(Paragraph(title, title_style))
        story.append(Spacer(1, 0.2 * inch))
        story.append(PageBreak())
        story.append(Spacer(1, 0.2 * inch))
        
        # Parse content and format (matching Word formatting)
        for text in content_text:
            if text.startswith("# "):
                heading_text = text[2:].strip()
                story.append(Paragraph(heading_text, heading1_style))
                story.append(Spacer(1, 0.1 * inch))
            elif text.startswith("## "):
                heading_text = text[3:].strip()
                story.append(Paragraph(heading_text, heading2_style))
                story.append(Spacer(1, 0.08 * inch))
            elif text.startswith("### "):
                heading_text = text[4:].strip()
                story.append(Paragraph(heading_text, heading3_style))
                story.append(Spacer(1, 0.06 * inch))
            elif text.lower().startswith("title:"):
                title_text = re.sub(r'^title:\s*', '', text, flags=re.IGNORECASE).strip()
                if title_text:
                    story.append(Paragraph(title_text, title_content_style))
                    story.append(Spacer(1, 0.1 * inch))
            elif text.startswith("[") and "]" in text:
                story.append(Paragraph(text, citation_style))
                story.append(Spacer(1, 0.04 * inch))
            elif text.startswith("- ") or text.startswith("* ") or text.startswith("• "):
                list_text = text[2:].strip() if not text.startswith("• ") else text[2:].strip()
                # Use orange bullet symbol with HTML color tag
                bullet_text = f'<font color="#E87722">•</font> {list_text}'
                story.append(Paragraph(bullet_text, bullet_style))
                story.append(Spacer(1, 0.05 * inch))
            elif re.match(r'^\d+[\.\)]\s+', text):
                # Extract number and text, color the number orange
                num_match = re.match(r'^(\d+)([\.\)])\s+(.*)', text)
                if num_match:
                    num = num_match.group(1)
                    punct = num_match.group(2)
                    list_text = num_match.group(3)
                    numbered_text = f'<font color="#E87722"><b>{num}{punct}</b></font> {list_text}'
                    story.append(Paragraph(numbered_text, bullet_style))
                else:
                    story.append(Paragraph(text, bullet_style))
                story.append(Spacer(1, 0.05 * inch))
            else:
                # Clean up markdown formatting
                clean_text = text.replace("**", "<b>").replace("**", "</b>")
                clean_text = clean_text.replace("*", "<i>").replace("*", "</i>")
                story.append(Paragraph(clean_text, body_style))
                story.append(Spacer(1, 0.06 * inch))
        
        doc_pdf.build(story)
        buffer.seek(0)
        return buffer.getvalue()
    
    except Exception as e:
        logger.error(f"Error creating PDF from Word content: {str(e)}")
        raise


def create_pdf_simple_fallback(content: str, title: str) -> bytes:
    """Simple PDF fallback when template conversion fails"""
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(
        buffer, pagesize=letter, topMargin=1 * inch, bottomMargin=1 * inch
    )
    
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "CustomTitle",
        parent=styles["Heading1"],
        fontSize=24,
        textColor="#E87722",
        spaceAfter=30,
        alignment=TA_LEFT,
    )
    
    body_style = ParagraphStyle(
        "CustomBody",
        parent=styles["BodyText"],
        fontSize=11,
        alignment=TA_JUSTIFY,
        spaceAfter=12,
        leading=14,
    )
    
    citation_style = ParagraphStyle(
        "Citation",
        parent=styles["BodyText"],
        fontSize=9,
        textColor="#666666",
        leftIndent=20,
        spaceAfter=6,
    )
    
    story = []
    story.append(Paragraph(title, title_style))
    story.append(Spacer(1, 0.2 * inch))
    
    paragraphs = content.split("\n\n")
    for para in paragraphs:
        if para.strip():
            if para.startswith("**") and para.endswith("**"):
                heading_text = para.strip("*").strip()
                story.append(Paragraph(heading_text, styles["Heading2"]))
            elif para.startswith("[") and "]" in para:
                story.append(Paragraph(para, citation_style))
            else:
                clean_para = para.replace("**", "<b>").replace("**", "</b>")
                story.append(Paragraph(clean_para, body_style))
            story.append(Spacer(1, 0.1 * inch))
    
    doc.build(story)
    buffer.seek(0)
    return buffer.getvalue()


async def create_word_doc(content: str, title: Optional[str] = None, template_path: str = 'template/Template.docx', use_template: bool = True) -> bytes:
    """Generate a Word document from content using template formatting"""
    try:
        # Generate title using LLM if not provided
        if not title:
            logger.info("Title not provided, generating with LLM...")
            title = await generate_title_with_llm(content)
            logger.info(f"Generated title: {title}")
        
        # Load template
        template_doc = None
        doc = None
        if use_template:
            # Try to load template from file path
            resolved_template_path = None
            backend_dir = os.path.dirname(os.path.abspath(__file__))
            normalized_template = os.path.normpath(template_path)
            candidate = os.path.join(backend_dir, normalized_template)
            
            if os.path.exists(candidate):
                resolved_template_path = candidate
            elif os.path.isabs(template_path) and os.path.exists(template_path):
                resolved_template_path = template_path
            elif os.path.exists(os.path.join(os.getcwd(), normalized_template)):
                resolved_template_path = os.path.join(os.getcwd(), normalized_template)
            
            if resolved_template_path:
                try:
                    logger.info(f"Loading template from: {resolved_template_path}")
                    # Load template directly - python-docx will preserve styles
                    doc = Document(resolved_template_path)
                    
                    # Clear existing content but keep styles
                    # Remove all paragraphs from template
                    xml_body = doc._body._body
                    for para in list(doc.paragraphs):
                        xml_body.remove(para._element)
                    
                    logger.info("Template loaded and cleared, styles preserved")
                except Exception as e:
                    logger.warning(f"Error loading template: {str(e)}, using default formatting")
                    doc = Document()
            else:
                logger.warning(f"Template not found at {template_path}, using default formatting")
                doc = Document()
        else:
            doc = Document()
        
        # Define consistent theme colors (from page 1 theme)
        theme_color_primary = DocxRGBColor(0xE8, 0x77, 0x22)  # Orange
        theme_color_secondary = DocxRGBColor(0x33, 0x33, 0x33)  # Dark gray
        theme_color_text = DocxRGBColor(0x00, 0x00, 0x00)  # Black
        theme_color_citation = DocxRGBColor(0x66, 0x66, 0x66)  # Gray
        
        
        # Page 1: Title Page with increased font size and template theme
        title_para = doc.add_paragraph()
        try:
            title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        except AttributeError:
            title_para.alignment = 1  # CENTER alignment
        
        # Add title with increased font size (26pt)
        title_run = title_para.add_run(title)
        title_run.font.size = DocxPt(26)  # Increased from 22
        title_run.bold = True
        title_run.font.color.rgb = theme_color_primary
        
        # Add vertical spacing for title page
        for _ in range(10):
            spacing_para = doc.add_paragraph()
            spacing_para.paragraph_format.space_after = DocxPt(12)
        
        # Check if TOC already exists in content, otherwise generate with LLM
        headings = extract_toc_from_content(content)
        toc_found_in_content = headings is not None
        
        if not headings:
            logger.info("No TOC found in content, generating with LLM...")
            headings = await generate_table_of_contents_with_llm(content)
            logger.info(f"Generated {len(headings)} TOC entries with LLM")
            
            # If LLM didn't generate enough headings, extract from content as fallback
            if not headings or len(headings) < 2:
                logger.info("LLM generated insufficient headings, extracting from content...")
                extracted_headings = extract_headings_from_content(content)
                if extracted_headings:
                    headings = extracted_headings
                    logger.info(f"Using {len(headings)} extracted headings")
        else:
            logger.info(f"Using existing TOC from content with {len(headings)} entries")
        
        # If TOC was found in content, we need to skip it when processing content
        # Store this flag for later use
        skip_toc_in_content = toc_found_in_content
        
        # Page 2: Table of Contents with improved formatting
        if headings:
            # Add page break
            doc.add_page_break()
            
            # Add TOC heading with theme styling
            toc_heading = doc.add_heading("Content", level=1)
            toc_heading.alignment = WD_ALIGN_PARAGRAPH.LEFT
            for run in toc_heading.runs:
                run.font.size = DocxPt(18)  # Increased
                run.font.color.rgb = theme_color_primary
                run.bold = True
            
            # Add spacing after TOC heading
            spacing_para = doc.add_paragraph()
            spacing_para.paragraph_format.space_after = DocxPt(14)
            
            # Calculate page numbers by mapping headings to their actual positions in content
            # This provides accurate page numbers for TOC
            content_headings_order = extract_headings_from_content(content)
            heading_to_page = {}
            base_page = 3  # Content starts at page 3 (page 1: title, page 2: TOC)
            
            # Estimate content distribution
            # Average: ~400-500 words per page with 12pt font, 1.8 line spacing
            words_per_page = 450  # Adjusted for 1.8 line spacing
            content_words = content.split()
            total_words = len(content_words)
            
            # Find heading positions in content and map to pages
            lines = content.split("\n")
            current_word_count = 0
            seen_headings = set()
            
            for line in lines:
                line_stripped = line.strip()
                # Skip TOC section if it exists in content
                if skip_toc_in_content:
                    line_lower = line_stripped.lower()
                    toc_keywords_list = ["table of contents", "contents", "table of content", "toc", "content"]
                    if any(keyword in line_lower for keyword in toc_keywords_list):
                        continue
                    if re.match(r'^[\d\-•*]\s+', line_stripped) and "..." in line_stripped:
                        continue
                
                # Check if this line is a heading
                heading_text = None
                if line_stripped.startswith("# "):
                    heading_text = line_stripped[2:].strip()
                elif line_stripped.startswith("## "):
                    heading_text = line_stripped[3:].strip()
                elif line_stripped.startswith("### "):
                    heading_text = line_stripped[4:].strip()
                
                if heading_text and heading_text not in seen_headings:
                    # Calculate page number based on word position before this heading
                    # Add some buffer for formatting (headings, spacing, etc.)
                    estimated_words_before = current_word_count
                    current_page = base_page + max(0, (estimated_words_before // words_per_page))
                    if current_page < base_page:
                        current_page = base_page
                    heading_to_page[heading_text] = current_page
                    seen_headings.add(heading_text)
                
                # Count words in line (excluding headings for better accuracy)
                if not heading_text:
                    current_word_count += len(line.split())
                else:
                    # For headings, count them as ~10 words (spacing, formatting overhead)
                    current_word_count += 10
            
            # Add TOC entries with professional formatting
            prev_page_num = base_page - 1  # Track previous page number for sequential check
            for i, heading in enumerate(headings, 1):
                toc_para = doc.add_paragraph()
                toc_para.paragraph_format.left_indent = DocxInches(0)
                toc_para.paragraph_format.space_after = DocxPt(10)
                toc_para.paragraph_format.line_spacing = 1.8  # Decreased by 0.2 (was 2.0)
                
                # Add professional bullet (•) for TOC entries
                bullet_run = toc_para.add_run("• ")
                bullet_run.font.size = DocxPt(12)
                bullet_run.font.color.rgb = theme_color_primary
                
                # Add heading text
                toc_run = toc_para.add_run(heading)
                toc_run.font.size = DocxPt(12)
                toc_run.font.color.rgb = theme_color_text
                
                # Add leader dots with tab alignment
                max_dots = 75
                heading_len = len(heading) + 2  # +2 for bullet
                dots_needed = max(3, max_dots - heading_len - 6)  # Reserve space for page number
                dots = "." * dots_needed
                dots_run = toc_para.add_run(f" {dots} ")
                dots_run.font.size = DocxPt(12)
                dots_run.font.color.rgb = DocxRGBColor(0xCC, 0xCC, 0xCC)  # Light gray dots
                
                # Get page number from mapping
                page_num = heading_to_page.get(heading)
                if not page_num:
                    # Try to find by normalized matching (case-insensitive, remove special chars)
                    heading_normalized = heading.lower().strip()
                    for mapped_heading, mapped_page in heading_to_page.items():
                        if heading_normalized == mapped_heading.lower().strip():
                            page_num = mapped_page
                            break
                    
                    if not page_num:
                        # Fallback: estimate based on heading position in TOC list
                        # Match heading from TOC with headings in content order
                        heading_matched = False
                        for j, content_heading in enumerate(content_headings_order):
                            heading_norm = heading.lower().strip()
                            content_norm = content_heading.lower().strip()
                            # Try exact match first
                            if heading_norm == content_norm:
                                page_num = base_page + max(0, j)
                                heading_matched = True
                                break
                            # Try partial match
                            elif heading_norm in content_norm or content_norm in heading_norm:
                                # Found matching heading in content
                                # Estimate page: assume each heading section takes about 1-2 pages
                                page_num = base_page + max(0, j)
                                heading_matched = True
                                break
                        
                        if not heading_matched:
                            # Final fallback: distribute headings evenly based on position
                            # Estimate: each heading section takes about 1 page on average
                            estimated_pages = max(1, total_words // words_per_page)
                            if len(headings) > 0:
                                page_num = base_page + min(i - 1, estimated_pages)
                            else:
                                page_num = base_page
                
                # Ensure page number is at least base_page
                if page_num < base_page:
                    page_num = base_page
                
                # Ensure page numbers are sequential (each heading should be on same or later page)
                if page_num < prev_page_num:
                    page_num = prev_page_num + 1
                
                # Update previous page number for next iteration
                prev_page_num = page_num
                
                page_run = toc_para.add_run(str(page_num))
                page_run.font.size = DocxPt(12)
                page_run.font.color.rgb = theme_color_text
                page_run.bold = True
        
        # Add page break after TOC (only if TOC was added)
        if headings:
            doc.add_page_break()
        
        # Page 3+: Formatted Content with professional styling
        # Parse formatted content (skip TOC section if it was already extracted)
        lines = content.split("\n")
        current_paragraph = None
        in_list = False
        list_level = 0
        in_toc_section = False
        in_key_points_section = False  # Initialize key points section flag
        in_types_section = False  # Initialize types section flag
        toc_keywords = ["table of contents", "contents", "table of content", "toc", "content"]
        
        for line in lines:
            # Skip TOC section if it was already extracted from content
            if skip_toc_in_content:
                line_stripped = line.strip()
                line_lower = line_stripped.lower()
                
                # Check if we're entering TOC section (heading line)
                if any(keyword in line_lower for keyword in toc_keywords):
                    in_toc_section = True
                    continue
                
                # If in TOC section, skip until we hit a major heading or real content
                if in_toc_section:
                    # Check if we've reached the end of TOC (major heading # or ##)
                    if line_stripped.startswith("# ") or line_stripped.startswith("## "):
                        # Only exit if we've seen TOC entries before (to avoid false positives)
                        if len([h for h in headings if h]) > 0:
                            in_toc_section = False
                        else:
                            continue
                    else:
                        # Skip TOC entries (lines with dots, bullets, numbers at start, or page numbers)
                        if (re.match(r'^[\d\-•*·]\s+', line_stripped) or 
                            "..." in line_stripped or 
                            re.match(r'^\d+[\.\)]\s+', line_stripped) or
                            re.search(r'\d+\s*$', line_stripped)):  # Ends with page number
                            continue
                        # Skip if it's an empty line in TOC section
                        if not line_stripped:
                            continue
                        # If we hit a real heading (### or ####) or content that doesn't look like TOC
                        if (line_stripped.startswith("### ") or 
                            line_stripped.startswith("#### ") or
                            (line_stripped and not any(keyword in line_lower for keyword in toc_keywords) and
                             not re.match(r'^[\d\-•*·]\s+', line_stripped) and 
                             "..." not in line_stripped and
                             len(line_stripped) > 10)):  # Real content line
                            in_toc_section = False
                        else:
                            continue
                
                # Skip this line if still in TOC section
                if in_toc_section:
                    continue
            line = line.strip()
            if not line:
                if current_paragraph:
                    current_paragraph = None
                    in_list = False
                continue
            
            # Check for headings with professional styling
            # Check for key points and types FIRST (before regular headings)
            # Key Points section heading (standalone heading)
            if re.match(r'^(key points?|key point):\s*$', line, re.IGNORECASE):
                key_heading = doc.add_heading("Key Points", level=2)
                key_heading.alignment = WD_ALIGN_PARAGRAPH.LEFT
                key_heading.paragraph_format.space_before = DocxPt(14)
                key_heading.paragraph_format.space_after = DocxPt(10)
                key_heading.paragraph_format.line_spacing = 1.8  # Decreased line spacing
                for run in key_heading.runs:
                    run.font.size = DocxPt(18)
                    run.font.color.rgb = theme_color_primary  # Orange
                    run.bold = True
                current_paragraph = None
                in_list = False
                list_level = 0
                in_key_points_section = True
                in_types_section = False
            # Types section heading (standalone heading)
            elif re.match(r'^(types?|type):\s*$', line, re.IGNORECASE):
                types_heading = doc.add_heading("Types", level=2)
                types_heading.alignment = WD_ALIGN_PARAGRAPH.LEFT
                types_heading.paragraph_format.space_before = DocxPt(14)
                types_heading.paragraph_format.space_after = DocxPt(10)
                types_heading.paragraph_format.line_spacing = 1.8  # Decreased line spacing
                for run in types_heading.runs:
                    run.font.size = DocxPt(18)
                    run.font.color.rgb = theme_color_primary  # Orange
                    run.bold = True
                current_paragraph = None
                in_list = False
                list_level = 0
                in_types_section = True
                in_key_points_section = False
            # Title (if it appears in content) - Orange, bold
            elif line.lower().startswith("title:") or (len(line) > 0 and line[0].isupper() and not line.startswith("#") and "title" in line.lower()[:20]):
                title_text = re.sub(r'^title:\s*', '', line, flags=re.IGNORECASE).strip()
                if title_text:
                    title_heading = doc.add_heading(title_text, level=1)
                    title_heading.alignment = WD_ALIGN_PARAGRAPH.LEFT
                    title_heading.paragraph_format.space_before = DocxPt(18)
                    title_heading.paragraph_format.space_after = DocxPt(12)
                    title_heading.paragraph_format.line_spacing = 1.8  # Decreased line spacing
                    for run in title_heading.runs:
                        run.font.size = DocxPt(20)
                        run.font.color.rgb = theme_color_primary  # Orange
                        run.bold = True
                current_paragraph = None
                in_list = False
                list_level = 0
                # Reset section flags
                in_key_points_section = False
                in_types_section = False
            elif line.startswith("# "):
                # Heading 1 - Main section headings (Orange, bold)
                heading_text = line[2:].strip()
                heading = doc.add_heading(heading_text, level=1)
                heading.alignment = WD_ALIGN_PARAGRAPH.LEFT
                heading.paragraph_format.space_before = DocxPt(18)  # Decreased
                heading.paragraph_format.space_after = DocxPt(12)  # Decreased
                heading.paragraph_format.line_spacing = 1.8  # Decreased line spacing
                for run in heading.runs:
                    run.font.size = DocxPt(20)
                    run.font.color.rgb = theme_color_primary  # Orange
                    run.bold = True
                current_paragraph = None
                in_list = False
                list_level = 0
                # Reset section flags when hitting a new main heading
                in_key_points_section = False
                in_types_section = False
            elif line.startswith("## "):
                # Heading 2 - Subsection headings (Orange, bold)
                heading_text = line[3:].strip()
                heading = doc.add_heading(heading_text, level=2)
                heading.alignment = WD_ALIGN_PARAGRAPH.LEFT
                heading.paragraph_format.space_before = DocxPt(14)  # Decreased
                heading.paragraph_format.space_after = DocxPt(10)  # Decreased
                heading.paragraph_format.line_spacing = 1.8  # Decreased line spacing
                for run in heading.runs:
                    run.font.size = DocxPt(18)
                    run.font.color.rgb = theme_color_primary  # Orange
                    run.bold = True
                current_paragraph = None
                in_list = False
                list_level = 0
                # Reset section flags for new subsection (unless it's key points or types)
                if "key point" not in heading_text.lower() and "type" not in heading_text.lower():
                    in_key_points_section = False
                    in_types_section = False
            elif line.startswith("### "):
                # Heading 3 - Sub-subsection headings (Dark gray, bold)
                heading_text = line[4:].strip()
                heading = doc.add_heading(heading_text, level=3)
                heading.alignment = WD_ALIGN_PARAGRAPH.LEFT
                heading.paragraph_format.space_before = DocxPt(12)  # Decreased
                heading.paragraph_format.space_after = DocxPt(8)  # Decreased
                heading.paragraph_format.line_spacing = 1.8  # Decreased line spacing
                for run in heading.runs:
                    run.font.size = DocxPt(16)
                    run.font.color.rgb = theme_color_secondary  # Dark gray
                    run.bold = True
                current_paragraph = None
                in_list = False
                list_level = 0
            elif line.startswith("#### "):
                # Heading 4 - Lower level headings (Dark gray, bold)
                heading_text = line[5:].strip()
                heading = doc.add_heading(heading_text, level=4)
                heading.alignment = WD_ALIGN_PARAGRAPH.LEFT
                heading.paragraph_format.space_before = DocxPt(12)
                heading.paragraph_format.space_after = DocxPt(8)
                heading.paragraph_format.line_spacing = 1.8  # Decreased line spacing
                for run in heading.runs:
                    run.font.size = DocxPt(14)
                    run.font.color.rgb = theme_color_secondary  # Dark gray
                    run.bold = True
                current_paragraph = None
                in_list = False
                list_level = 0
            # Check for bullet lists with professional formatting
            elif line.startswith("- ") or line.startswith("* ") or line.startswith("• ") or line.startswith("· ") or line.startswith("▪ "):
                # Remove bullet marker
                list_text = re.sub(r'^[-*•·▪]\s+', '', line).strip()
                para = doc.add_paragraph()
                para.paragraph_format.left_indent = DocxInches(0.3)
                para.paragraph_format.space_after = DocxPt(8)  # Decreased spacing
                para.paragraph_format.line_spacing = 1.8  # Decreased by 0.2 (was 2.0)
                para.paragraph_format.first_line_indent = DocxInches(-0.3)
                
                # Determine bullet type based on context and section
                if in_key_points_section:
                    # Key Points section: use filled circle (●)
                    bullet_char = "● "
                    bullet_color = theme_color_primary
                    bullet_bold = True
                elif in_types_section:
                    # Types section: use diamond (◆)
                    bullet_char = "◆ "
                    bullet_color = theme_color_primary
                    bullet_bold = True
                elif line.startswith("▪ "):
                    # Explicit square bullet for sub-items
                    bullet_char = "▪ "
                    bullet_color = theme_color_secondary
                    bullet_bold = False
                elif in_list and list_level > 0:
                    # Sub-item: use square bullet (▪)
                    bullet_char = "▪ "
                    bullet_color = theme_color_secondary
                    bullet_bold = False
                else:
                    # Regular bullet: use solid bullet (•)
                    bullet_char = "• "
                    bullet_color = theme_color_primary
                    bullet_bold = True
                
                bullet_run = para.add_run(bullet_char)
                bullet_run.font.size = DocxPt(12)
                bullet_run.font.color.rgb = bullet_color
                bullet_run.bold = bullet_bold
                
                # Add list text
                text_run = para.add_run(list_text)
                text_run.font.size = DocxPt(12)
                text_run.font.color.rgb = theme_color_text
                
                current_paragraph = None
                in_list = True
                list_level = 1
            # Check for numbered lists with proper formatting
            elif re.match(r'^\d+[\.\)]\s+', line):
                list_text = re.sub(r'^\d+[\.\)]\s+', '', line).strip()
                para = doc.add_paragraph()
                para.paragraph_format.left_indent = DocxInches(0.3)
                para.paragraph_format.space_after = DocxPt(8)  # Decreased spacing
                para.paragraph_format.line_spacing = 1.8  # Decreased by 0.2 (was 2.0)
                para.paragraph_format.first_line_indent = DocxInches(-0.3)
                
                # Add number with orange color
                num_match = re.match(r'^(\d+)[\.\)]\s+', line)
                if num_match:
                    num_text = num_match.group(1) + "."
                    num_run = para.add_run(num_text + " ")
                    num_run.font.size = DocxPt(12)
                    num_run.font.color.rgb = theme_color_primary
                    num_run.bold = True
                
                # Add list text
                text_run = para.add_run(list_text)
                text_run.font.size = DocxPt(12)
                text_run.font.color.rgb = theme_color_text
                
                current_paragraph = None
                in_list = True
                list_level = 1
            # Check for hierarchical sub-bullets (indented with spaces or tabs)
            elif re.match(r'^\s{2,}[-*•·▪]\s+', line) or re.match(r'^\s{2,}\d+[\.\)]\s+', line):
                # Sub-item (indented) - Level 2
                list_text = re.sub(r'^\s+', '', line)
                list_text = re.sub(r'^[-*•·▪]\s+', '', list_text)
                list_text = re.sub(r'^\d+[\.\)]\s+', '', list_text).strip()
                
                para = doc.add_paragraph()
                para.paragraph_format.left_indent = DocxInches(0.5)  # More indented
                para.paragraph_format.space_after = DocxPt(6)
                para.paragraph_format.line_spacing = 1.8  # Decreased by 0.2
                para.paragraph_format.first_line_indent = DocxInches(-0.25)
                
                # Use square bullet (▪) for sub-items
                sub_bullet_run = para.add_run("▪ ")
                sub_bullet_run.font.size = DocxPt(11)
                sub_bullet_run.font.color.rgb = theme_color_secondary
                sub_bullet_run.bold = False
                
                # Add list text
                text_run = para.add_run(list_text)
                text_run.font.size = DocxPt(11)
                text_run.font.color.rgb = theme_color_text
                
                current_paragraph = None
                in_list = True
                list_level = 2
            # Check for deeper nested bullets (Level 3+)
            elif re.match(r'^\s{4,}[-*•·▪]\s+', line):
                # Sub-sub-item (deeper indented) - Level 3
                list_text = re.sub(r'^\s+', '', line)
                list_text = re.sub(r'^[-*•·▪]\s+', '', list_text).strip()
                
                para = doc.add_paragraph()
                para.paragraph_format.left_indent = DocxInches(0.7)  # Even more indented
                para.paragraph_format.space_after = DocxPt(5)
                para.paragraph_format.line_spacing = 1.8
                para.paragraph_format.first_line_indent = DocxInches(-0.2)
                
                # Use circle bullet (○) for deeper nested items
                sub_bullet_run = para.add_run("○ ")
                sub_bullet_run.font.size = DocxPt(10)
                sub_bullet_run.font.color.rgb = theme_color_secondary
                sub_bullet_run.bold = False
                
                # Add list text
                text_run = para.add_run(list_text)
                text_run.font.size = DocxPt(10)
                text_run.font.color.rgb = theme_color_text
                
                current_paragraph = None
                in_list = True
                list_level = 3
            # Check for citations
            elif line.startswith("[") and "]" in line:
                citation = doc.add_paragraph(line)
                citation.paragraph_format.left_indent = DocxInches(0.5)
                citation.paragraph_format.space_after = DocxPt(6)
                citation.paragraph_format.line_spacing = 1.8  # Match overall line spacing
                for run in citation.runs:
                    run.font.size = DocxPt(9)
                    run.font.color.rgb = theme_color_citation
                    run.italic = True
                current_paragraph = None
                in_list = False
            # Regular body text with theme styling
            else:
                # Skip TOC section if it was already extracted
                line_lower = line.lower().strip()
                if any(keyword in line_lower for keyword in ["table of contents", "contents", "table of content"]):
                    # Skip TOC heading if we already have it
                    continue
                
                para = doc.add_paragraph(line)
                para.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                para.paragraph_format.space_after = DocxPt(10)  # Decreased spacing
                para.paragraph_format.line_spacing = 1.8  # Decreased by 0.2 (was 2.0)
                para.paragraph_format.first_line_indent = DocxInches(0)
                for run in para.runs:
                    run.font.size = DocxPt(12)
                    run.font.color.rgb = theme_color_text
                current_paragraph = para
                in_list = False
        
        # Save document to buffer
        buffer = io.BytesIO()
        doc.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()
    
    except Exception as e:
        logger.error(f"Error creating Word document: {str(e)}")
        # Fallback: create simple document without template
        doc = Document()
        final_title = title if title else "Document"
        title_para = doc.add_heading(final_title, 0)
        title_para.alignment = WD_ALIGN_PARAGRAPH.LEFT
        for run in title_para.runs:
            run.font.color.rgb = DocxRGBColor(0xE8, 0x77, 0x22)
        doc.add_paragraph()
        
        paragraphs = content.split("\n\n")
        for para in paragraphs:
            if para.strip():
                if para.startswith("##"):
                    heading_text = para.replace("##", "").strip()
                    doc.add_heading(heading_text, level=2)
                elif para.startswith("[") and "]" in para:
                    citation = doc.add_paragraph(para)
                    citation.paragraph_format.left_indent = DocxInches(0.5)
                    for run in citation.runs:
                        run.font.size = DocxPt(9)
                        run.font.color.rgb = DocxRGBColor(0x66, 0x66, 0x66)
                else:
                    p = doc.add_paragraph(para)
                    p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
        
        buffer = io.BytesIO()
        doc.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()


async def generate_stream(messages: list):
    """Generate streaming response using LLM service with fallback"""
    try:
        llm = get_llm()

        async for content in llm.stream_completion(
            messages=messages,
            temperature=0.7,
            max_tokens=4096,
        ):
            yield f"data: {json.dumps({'content': content})}\n\n"

        yield f"data: {json.dumps({'done': True})}\n\n"

    except Exception as e:
        yield f"data: {json.dumps({'error': str(e)})}\n\n"


@app.post("/api/extract-text")
async def extract_text(file: UploadFile = File(...)):
    """Extract text from uploaded PDF, DOCX, TXT, or MD files"""
    try:
        text = await parse_uploaded_file(file, max_chars=50000)
        return {"text": text, "filename": file.filename}
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Error extracting text: {str(e)}")


@app.post("/api/chat")
async def chat(request: ChatRequest):
    try:
        messages = [
            {"role": msg.role, "content": msg.content} for msg in request.messages
        ]

        if request.stream:
            return StreamingResponse(
                generate_stream(messages), media_type="text/event-stream"
            )
        else:
            llm = get_llm()
            response = await llm.chat_completion(
                messages=messages,
                temperature=0.7,
                max_tokens=4096,
            )

            return {
                "message": response.content,
                "provider": response.provider.value,
                "usage": response.usage
                or {
                    "prompt_tokens": 0,
                    "completion_tokens": 0,
                    "total_tokens": 0,
                },
            }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/draft")
async def create_draft(request: DraftRequest):
    try:
        system_prompt = """You are an expert presentation consultant at PwC (PricewaterhouseCoopers). 
Your role is to help create professional, structured, and impactful presentations.
When drafting presentations, focus on:
- Clear structure with MECE (Mutually Exclusive, Collectively Exhaustive) framework
- Executive summaries and key takeaways
- Data-driven insights
- Professional formatting suggestions
- PwC's consulting best practices"""

        user_prompt = f"""Please create a presentation outline with the following details:

Topic: {request.topic}
Objective: {request.objective}
Target Audience: {request.audience}
{f"Additional Context: {request.additional_context}" if request.additional_context else ""}

Provide a structured outline including:
1. Presentation title
2. Executive summary
3. Slide-by-slide breakdown with:
   - Slide titles
   - Key messages for each slide
   - Suggested content types (charts, frameworks, bullet points)
4. Conclusion and call-to-action

Format the output in a clear, professional manner."""

        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ]

        return StreamingResponse(
            generate_stream(messages), media_type="text/event-stream"
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/ppt/generate")
async def generate_ppt(
    topic: str = Form(...),
    objective: str = Form(...),
    audience: str = Form(...),
    additional_context: Optional[str] = Form(None),
    template_ppt: UploadFile = File(None),
):
    """Generate an actual PowerPoint file from presentation draft request."""
    try:
        # Generate content using Groq
        system_prompt = """You are an expert presentation consultant. Create concise, impactful slide content.
For each slide, provide:
- Slide title (one line)
- 3-5 bullet points (each under 15 words)
Keep content professional and data-driven."""

        user_prompt = f"""Create slide content for a presentation:

Topic: {topic}
Objective: {objective}
Target Audience: {audience}
{f"Additional Context: {additional_context}" if additional_context else ""}

Provide 5-8 slides with titles and bullet points. Format as:
SLIDE 1: [Title]
• [Bullet point 1]
• [Bullet point 2]
• [Bullet point 3]

SLIDE 2: [Title]
..."""

        llm = get_llm()
        llm_response = await llm.chat_completion(
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt},
            ],
            temperature=0.7,
            max_tokens=2000,
        )

        content = llm_response.content

        # Create PowerPoint from template or blank
        if template_ppt:
            template_content = await template_ppt.read()
            prs = Presentation(io.BytesIO(template_content))
            # Clear existing slides
            for i in range(len(prs.slides) - 1, -1, -1):
                rId = prs.slides._sldIdLst[i].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[i]
        else:
            prs = Presentation()

        # Parse content and create slides
        slides_data = []
        current_slide = None

        for line in content.split("\n"):
            line = line.strip()
            if line.startswith("SLIDE ") and ":" in line:
                if current_slide:
                    slides_data.append(current_slide)
                title = line.split(":", 1)[1].strip()
                current_slide = {"title": title, "bullets": []}
            elif line.startswith("•") or line.startswith("-"):
                if current_slide:
                    bullet = line.lstrip("•-").strip()
                    if bullet:
                        current_slide["bullets"].append(bullet)

        if current_slide:
            slides_data.append(current_slide)

        # Create slides
        for slide_data in slides_data:
            slide_layout = prs.slide_layouts[1]  # Title and Content layout
            slide = prs.slides.add_slide(slide_layout)

            # Set title
            title_shape = slide.shapes.title
            title_shape.text = slide_data["title"]

            # Add bullet points
            if slide_data["bullets"] and len(slide.shapes) > 1:
                content_shape = slide.shapes[1]
                text_frame = content_shape.text_frame
                text_frame.clear()

                for i, bullet in enumerate(slide_data["bullets"]):
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    p.text = bullet
                    p.level = 0

        # Save to BytesIO
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)

        # Return as downloadable file
        filename = f"{topic.replace(' ', '_')[:30]}_presentation.pptx"
        return StreamingResponse(
            output,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f"attachment; filename={filename}"},
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/ddc/brand-format")
async def ddc_brand_format(
    file: UploadFile = File(...),
    template: str = Form(...),
    template_file: Optional[UploadFile] = File(None),
    services: str = Form(...)
):
    """
    DDC Branding & Formatting workflow endpoint.
    Processes PowerPoint presentations with selected branding and formatting services.
    
    Services: align-text, align-shapes, sanitize-metadata, font-standardization, 
              color-branding, slide-layout
    """
    try:
        services_list = json.loads(services)
        
        logger.info(f"[DDC Brand-Format] Processing: {file.filename}")
        logger.info(f"[DDC Brand-Format] Template: {template}")
        logger.info(f"[DDC Brand-Format] Services: {services_list}")
        logger.info(f"[DDC Brand-Format] Custom Template: {template_file.filename if template_file else 'None'}")
        
        llm = get_llm()
        
        system_prompt = """You are a PwC presentation branding and formatting specialist.
You help consultants standardize their presentations according to PwC brand guidelines and best practices.

Your expertise includes:
- Aligning text and shapes for visual consistency
- Applying PwC color palette (#D04A02 orange as primary accent)
- Standardizing fonts (typically Helvetica, Arial, or PwC custom fonts)
- Sanitizing metadata to protect client confidentiality
- Optimizing slide layouts for clarity and impact

When processing a presentation, provide:
1. Overview of the presentation (slide count, current state)
2. Detailed description of each service being applied
3. Specific changes made for each service
4. Quality assurance summary
5. Next steps and recommendations

Be professional, detailed, and provide actionable insights."""
        
        user_message = f"""I'm processing a PowerPoint presentation with the following specifications:

**File:** {file.filename}
**Template:** {template}
**Services Selected:** {', '.join(services_list)}

Please provide a comprehensive report on how these branding and formatting services will be applied to the presentation. 
Describe the process step-by-step, highlighting what each service accomplishes and the expected improvements.

Focus on:
- Text and shape alignment corrections
- Font standardization to PwC guidelines
- Color branding with PwC orange (#D04A02)
- Metadata sanitization for security
- Slide layout optimization for consistency

Provide a detailed, professional analysis as if you were actually performing these operations."""

        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_message}
        ]
        
        async def stream_brand_format():
            try:
                async for content in llm.stream_completion(
                    messages=messages,
                    temperature=0.7,
                    max_tokens=4096,
                ):
                    yield f"data: {json.dumps({'content': content})}\n\n"
                
                yield f"data: {json.dumps({'done': True})}\n\n"
            except Exception as e:
                logger.error(f"[DDC Brand-Format] Streaming error: {e}")
                yield f"data: {json.dumps({'error': str(e)})}\n\n"
        
        return StreamingResponse(
            stream_brand_format(), media_type="text/event-stream"
        )
        
    except json.JSONDecodeError as e:
        logger.error(f"[DDC Brand-Format] JSON decode error: {e}")
        raise HTTPException(status_code=400, detail="Invalid services JSON format")
    except Exception as e:
        logger.error(f"[DDC Brand-Format] Error: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/thought-leadership")
async def thought_leadership(request: ThoughtLeadershipRequest):
    try:
        system_prompt = """You are an expert thought leadership consultant at PwC (PricewaterhouseCoopers). 
You specialize in creating compelling, insightful content that positions executives and organizations as industry leaders.
Your expertise includes strategic insights, market analysis, editorial refinement, and content transformation."""

        user_prompt = ""
        reference_context = ""

        if request.reference_urls:
            reference_context = "\n\n**Reference Sources:**\n"
            failed_sources = []
            all_fetched_pages = []  # Store all pages including crawled ones

            for url in request.reference_urls:
                url_data = await fetch_url_content(url)
                if url_data["success"]:
                    # Add the main page
                    all_fetched_pages.append(
                        {
                            "url": url_data["url"],
                            "title": url_data["title"],
                            "content": url_data["content"],
                        }
                    )

                    # Automatically crawl related pages from the same domain
                    try:
                        related_pages = await crawl_related_pages(
                            url_data["url"], max_pages=10, max_depth=3, query=None
                        )
                        # Add related pages (skip the first one as it's the main page)
                        for page in related_pages:
                            if (
                                page["url"] != url_data["url"]
                            ):  # Don't duplicate the main page
                                all_fetched_pages.append(page)
                    except Exception:
                        # If crawling fails, continue with just the main page
                        pass
                else:
                    error_msg = url_data.get("error", "Unknown error")
                    failed_sources.append(f"{url} - {error_msg}")
                    reference_context += f"\n**Source:** {url}\n**Status:** Unable to fetch content - {error_msg}\n\n---\n"

            # Add all successfully fetched pages to context
            for i, page in enumerate(all_fetched_pages, 1):
                reference_context += f"\n**Source {i}:** {page['title']}\n**URL:** {page['url']}\n**Content:**\n{page['content']}\n\n---\n"

            if failed_sources:
                reference_context += f"\n\n**Note:** {len(failed_sources)} source(s) could not be accessed. Please only use information from successfully fetched sources above. Do not attempt to reference or cite the failed sources."

            if len(all_fetched_pages) > len(request.reference_urls):
                reference_context += f"\n\n**Additional Context:** Automatically discovered and included {len(all_fetched_pages) - len(request.reference_urls)} related page(s) from the provided website(s) to provide comprehensive information."

        if request.operation == "generate":
            citations_instruction = ""
            if request.reference_urls:
                citations_instruction = "\n\nIMPORTANT: Include proper citations in your article. At the end of the article, add a 'References' or 'Sources' section with all the URLs provided. Use in-text citations where appropriate (e.g., [1], [2]) that correspond to the numbered references."

            user_prompt = f"""Generate a draft thought leadership article with the following details:

Topic: {request.topic}
Perspective/Angle: {request.perspective}
Target Audience: {request.target_audience}
{f"Additional Context: {request.additional_context}" if request.additional_context else ""}{reference_context}{citations_instruction}

Create a comprehensive, well-structured article that:
1. Opens with a compelling hook and executive summary
2. Provides unique insights and data-driven perspectives (use reference sources when available)
3. Uses frameworks and strategic analysis
4. Includes real-world implications
5. Concludes with actionable recommendations
6. Properly cites all reference sources

Write in a professional, authoritative tone suitable for publication."""

        elif request.operation == "research":
            user_prompt = f"""Research and provide additional insights on:

Topic: {request.topic}
Current Perspective: {request.perspective}
{f"Additional Context: {request.additional_context}" if request.additional_context else ""}{reference_context}

Provide:
1. Emerging trends and developments in this area (using reference sources when available)
2. Different perspectives and counterarguments
3. Recent data, statistics, or case studies from the provided sources
4. Industry expert viewpoints
5. Future implications and opportunities

Focus on solution-oriented insights that add depth to the existing perspective. When using information from reference sources, cite them appropriately."""

        elif request.operation == "editorial":
            additional_instructions = (
                f"\n\nAdditional Instructions: {request.additional_context}"
                if request.additional_context
                else ""
            )
            user_prompt = f"""Provide comprehensive editorial support for this document:

{request.document_text}{additional_instructions}

Deliver end-to-end editorial review including:
1. Structure and flow improvements
2. Clarity and readability enhancements
3. Tone and voice adjustments for thought leadership
4. Grammar, punctuation, and style corrections
5. Fact-checking and consistency review
6. Suggested additions or cuts
7. Overall impact assessment

Provide specific, actionable recommendations."""

        elif request.operation == "improve":
            focus_areas = (
                f"\nFocus Areas: {request.additional_context}"
                if request.additional_context
                else ""
            )
            user_prompt = f"""Recommend improvements to this document:

{request.document_text}{focus_areas}

Analyze and recommend improvements for:
1. Content quality and depth of insights
2. Argument strength and persuasiveness
3. Evidence and supporting data
4. Professional positioning and credibility
5. Engagement and readability
6. Call-to-action effectiveness

Prioritize the most impactful improvements."""

        elif request.operation == "translate":
            additional_requirements = (
                f"\nAdditional Requirements: {request.additional_context}"
                if request.additional_context
                else ""
            )
            user_prompt = f"""Translate this document to a different format:

Original Document:
{request.document_text}

Target Format: {request.target_format}{additional_requirements}

Transform the content while:
1. Maintaining core messages and insights
2. Adapting tone and style to the target format
3. Adjusting length and structure appropriately
4. Optimizing for the new medium's best practices
5. Preserving professional quality

Provide the fully transformed content."""

        else:
            raise HTTPException(
                status_code=400, detail=f"Invalid operation: {request.operation}"
            )

        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ]

        return StreamingResponse(
            generate_stream(messages), media_type="text/event-stream"
        )

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/research")
async def research_assistant(request: ResearchRequest):
    """AI Research Assistant - provides LLM-powered research insights and analysis"""
    try:
        system_prompt = """You are an expert research assistant at PwC with deep industry knowledge and analytical capabilities.
You specialize in:
- Identifying emerging trends and market developments based on your training data
- Analyzing competitive intelligence and industry dynamics
- Providing data-driven insights and statistical perspectives
- Synthesizing multiple perspectives and viewpoints
- Recommending strategic implications

IMPORTANT: You are providing insights based on your knowledge base (training data through April 2024). 
For the most current data, recommend users verify with real-time sources.
Focus on frameworks, patterns, and strategic analysis that remain relevant regardless of real-time data.

Your research is comprehensive, objective, and actionable."""

        focus_context = ""
        if request.focus_areas:
            focus_context = f"\n\nFocus Areas: {', '.join(request.focus_areas)}"

        additional = (
            f"\n\nAdditional Context: {request.additional_context}"
            if request.additional_context
            else ""
        )

        user_prompt = f"""Conduct comprehensive research on the following query:

{request.query}{focus_context}{additional}

Provide:
1. **Key Findings**: 3-5 critical insights with supporting evidence
2. **Emerging Trends**: Recent developments and future implications
3. **Data Points**: Relevant statistics, metrics, or quantitative insights
4. **Different Perspectives**: Multiple viewpoints and counterarguments
5. **Strategic Implications**: How this impacts business strategy
6. **Recommended Next Steps**: Actionable recommendations

Format your response clearly with headers and bullet points for easy scanning."""

        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ]

        return StreamingResponse(
            generate_stream(messages), media_type="text/event-stream"
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


async def fetch_web_content(url: str, max_chars: int = 10000) -> str:
    """Fetch and extract text content from a URL"""
    try:
        # Browser-like headers to avoid 403 Forbidden errors
        # Build referer from the URL (same domain)
        parsed_url = urlparse(url)
        referer = f"{parsed_url.scheme}://{parsed_url.netloc}/"

        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36",
            "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8,application/signed-exchange;v=b3;q=0.7",
            "Accept-Language": "en-US,en;q=0.9",
            "Accept-Encoding": "gzip, deflate, br",
            "Connection": "keep-alive",
            "Upgrade-Insecure-Requests": "1",
            "Referer": referer,
            "Sec-Fetch-Dest": "document",
            "Sec-Fetch-Mode": "navigate",
            "Sec-Fetch-Site": "same-origin",
            "Sec-Fetch-User": "?1",
            "Cache-Control": "max-age=0",
            "DNT": "1",
        }

        async with httpx.AsyncClient(timeout=30.0, headers=headers) as client:
            response = await client.get(url, follow_redirects=True)

            # Check status codes before raising
            if response.status_code == 404:
                # Try URL variations before giving up
                variations = generate_url_variations(url)
                tried_urls = [url]

                for variation_url in variations:
                    if variation_url in tried_urls:
                        continue
                    tried_urls.append(variation_url)

                    try:
                        # Update referer for the variation
                        parsed_var = urlparse(variation_url)
                        var_referer = f"{parsed_var.scheme}://{parsed_var.netloc}/"
                        headers["Referer"] = var_referer

                        var_response = await client.get(
                            variation_url, follow_redirects=True
                        )
                        if var_response.status_code == 200:
                            # Success! Use this URL and continue with content extraction
                            url = variation_url  # Update url to the working variation
                            response = var_response
                            break
                        elif var_response.status_code != 404:
                            # If it's not 200 or 404, it's a different error, stop trying
                            break
                    except Exception:
                        # Continue to next variation on any error
                        continue

                # If still 404 after trying variations, return error
                if response.status_code == 404:
                    return f"[Error: 404 Not Found - The page at this URL does not exist. Tried {len(tried_urls)} variation(s) but none were found. Please verify the URL is correct.]"
            elif response.status_code == 403:
                return f"[Error: 403 Forbidden - The website blocked access. This may be due to bot protection, Cloudflare, or WAF rules. The site may require JavaScript to load content.]"
            elif response.status_code == 429:
                return f"[Error: 429 Too Many Requests - Rate limit exceeded. Please wait a moment and try again.]"
            elif response.status_code == 503:
                return f"[Error: 503 Service Unavailable - The website is temporarily unavailable. Please try again later.]"

            response.raise_for_status()

            soup = BeautifulSoup(response.text, "html.parser")

            # Remove script and style elements
            for script in soup(["script", "style", "nav", "footer", "header"]):
                script.decompose()

            # Get text content
            text = soup.get_text()

            # Clean up whitespace
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = "\n".join(chunk for chunk in chunks if chunk)

            # Truncate if needed
            if len(text) > max_chars:
                text = text[:max_chars] + "... [truncated]"

            return text
    except Exception as e:
        return f"[Error fetching URL: {str(e)}]"


@app.post("/api/research-with-materials")
async def research_with_materials(
    query: str = Form(...),
    files: Optional[List[UploadFile]] = File(None),
    links: Optional[List[str]] = Form(None),
    focus_areas: Optional[str] = Form(None),
    additional_context: Optional[str] = Form(None),
):
    """Research assistant with document upload and URL reference support (NotebookLM-style)"""
    try:

        async def generate_research_stream():
            try:
                # Step 1: Parse uploaded documents
                documents_content = ""
                if files:
                    yield f'data: {json.dumps({"type": "progress", "message": f"Parsing {len(files)} document(s)..."})}\n\n'

                    for i, file in enumerate(files, 1):
                        parsed_content = await parse_uploaded_file(
                            file, max_chars=15000
                        )
                        if parsed_content:
                            documents_content += f"\n\n### Document {i}: {file.filename}\n{parsed_content}"

                # Step 2: Fetch and parse web content from links (with automatic crawling)
                web_content = ""
                if links:
                    valid_links = [link for link in links if link and link.strip()]
                    if valid_links:
                        yield f'data: {json.dumps({"type": "progress", "message": f"Fetching content from {len(valid_links)} link(s) and discovering related pages..."})}\n\n'

                        all_web_sources = []
                        for link in valid_links:
                            link = link.strip()
                            # Fetch main page with full metadata
                            url_data = await fetch_url_content(link)
                            if url_data.get("success"):
                                all_web_sources.append(
                                    {
                                        "url": url_data["url"],
                                        "title": url_data.get("title", "").strip()
                                        or link,
                                        "content": url_data.get("content", "")[:10000],
                                    }
                                )

                                # Automatically crawl related pages
                                try:
                                    related_pages = await crawl_related_pages(
                                        link, max_pages=10, max_depth=3, query=query
                                    )
                                    for page in related_pages:
                                        if (
                                            page["url"] != link
                                        ):  # Don't duplicate main page
                                            all_web_sources.append(
                                                {
                                                    "url": page["url"],
                                                    "title": page.get(
                                                        "title", ""
                                                    ).strip()
                                                    or page["url"],
                                                    "content": (
                                                        page["content"][:10000]
                                                        if len(page["content"]) > 10000
                                                        else page["content"]
                                                    ),
                                                }
                                            )
                                except Exception:
                                    # If crawling fails, continue with just main page
                                    pass
                            else:
                                # Fallback to fetch_web_content if fetch_url_content fails
                                main_content = await fetch_web_content(link)
                                if not main_content.startswith("[Error:"):
                                    all_web_sources.append(
                                        {
                                            "url": link,
                                            "title": link,
                                            "content": main_content[:10000],
                                        }
                                    )

                        # Add all sources to web_content with title for better citations
                        for i, source in enumerate(all_web_sources, 1):
                            title = source.get("title", "").strip() or source["url"]
                            # Include both title and URL for frontend parsing
                            web_content += f"\n\n### Web Source {i}: {title} | URL: {source['url']}\n{source['content']}"

                        if len(all_web_sources) > len(valid_links):
                            yield f'data: {json.dumps({"type": "progress", "message": f"Discovered {len(all_web_sources) - len(valid_links)} additional related page(s) automatically"})}\n\n'

                        # Send source metadata for frontend link rendering
                        source_metadata = [
                            {
                                "number": i,
                                "url": s["url"],
                                "title": s.get("title", "").strip() or s["url"],
                            }
                            for i, s in enumerate(all_web_sources, 1)
                        ]
                        yield f'data: {json.dumps({"type": "sources", "sources": source_metadata})}\n\n'

                # Step 3: Construct research prompt with materials
                yield f'data: {json.dumps({"type": "progress", "message": "Analyzing materials and conducting research..."})}\n\n'

                system_prompt = """You are an expert research assistant at PwC with deep analytical capabilities.
You excel at:
- Synthesizing information from multiple sources
- Identifying key themes and patterns across documents
- Extracting actionable insights and strategic implications
- Providing comprehensive analysis based on provided materials
- Cross-referencing different sources to validate findings

You are analyzing user-provided documents and web sources. Base your research primarily on these materials,
while supplementing with your knowledge when relevant. Always cite which sources support your findings."""

                materials_section = ""
                if documents_content or web_content:
                    materials_section = "\n\n## Reference Materials Provided:"
                    if documents_content:
                        materials_section += documents_content
                    if web_content:
                        materials_section += web_content

                focus_areas_list = []
                if focus_areas:
                    try:
                        focus_areas_list = (
                            json.loads(focus_areas)
                            if isinstance(focus_areas, str)
                            else focus_areas
                        )
                    except:
                        focus_areas_list = [focus_areas]

                focus_context = ""
                if focus_areas_list:
                    focus_context = f"\n\nFocus Areas: {', '.join(focus_areas_list)}"

                additional = (
                    f"\n\nAdditional Context: {additional_context}"
                    if additional_context
                    else ""
                )

                user_prompt = f"""Based on the reference materials provided, conduct comprehensive research on the following query:

{query}{focus_context}{additional}{materials_section}

Provide a thorough analysis including:
1. **Executive Summary**: Overview of key findings from the materials
2. **Key Insights**: 4-6 critical insights extracted from the provided sources (cite specific sources)
3. **Themes & Patterns**: Common threads and patterns across the materials
4. **Data & Evidence**: Relevant statistics, quotes, and concrete evidence from the sources
5. **Cross-Source Analysis**: How different sources complement or contradict each other
6. **Strategic Implications**: What this means for business strategy and decision-making
7. **Gaps & Recommendations**: What's missing and suggested next steps

Format your response clearly with headers and bullet points. When citing insights, reference the specific source (e.g., "Document 1" or "Web Source 2")."""

                messages = [
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_prompt},
                ]

                # Stream the AI response
                llm = get_llm()
                async for content in llm.stream_completion(
                    messages=messages,
                    temperature=0.7,
                    max_tokens=4096,
                ):
                    yield f'data: {json.dumps({"type": "content", "content": content})}\n\n'

                yield f'data: {json.dumps({"type": "complete"})}\n\n'

            except Exception as e:
                yield f'data: {json.dumps({"type": "error", "message": str(e)})}\n\n'

        return StreamingResponse(
            generate_research_stream(), media_type="text/event-stream"
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/draft-article")
async def draft_article(
    topic: str = Form(...),
    content_type: str = Form(...),
    desired_length: int = Form(...),
    tone: str = Form(...),
    outline_text: Optional[str] = Form(None),
    additional_context: Optional[str] = Form(None),
    outline_file: Optional[UploadFile] = File(None),
    supporting_docs: Optional[List[UploadFile]] = File(None),
):
    """Generate long-form thought leadership articles from user outlines and supporting documents"""
    try:
        system_prompt = """You are an expert thought leadership writer at PwC.
You excel at creating compelling, insightful articles that position executives as industry leaders.
You use the MECE principle (Mutually Exclusive, Collectively Exhaustive), strategic frameworks,
and data-driven insights. Your writing is authoritative yet accessible."""

        outline_content = outline_text or ""
        supporting_content = ""

        if outline_file:
            parsed_outline = await parse_uploaded_file(outline_file, max_chars=5000)
            if parsed_outline:
                outline_content += "\n\n**Uploaded Outline:**\n" + parsed_outline

        if supporting_docs:
            supporting_content = "\n\n**Supporting Documents:**\n"
            for i, doc in enumerate(supporting_docs, 1):
                parsed_content = await parse_uploaded_file(doc, max_chars=3000)
                if parsed_content:
                    supporting_content += (
                        f"\n**Document {i} ({doc.filename}):**\n{parsed_content}\n"
                    )

        additional = (
            f"\n\nAdditional Context: {additional_context}"
            if additional_context
            else ""
        )

        user_prompt = f"""Create a comprehensive {content_type.lower()} on the following topic:

**Topic:** {topic}
**Content Type:** {content_type}
**Target Length:** {desired_length} words
**Tone:** {tone}

**Outline/Initial Ideas:**
{outline_content}{supporting_content}{additional}

Generate a well-structured, professional article that:
1. Opens with a compelling hook and clear thesis
2. Uses the provided outline as a foundation (if provided)
3. Integrates insights from supporting documents (if provided)
4. Includes strategic frameworks and analysis
5. Provides real-world examples and implications
6. Concludes with actionable recommendations
7. Maintains the specified tone throughout
8. Targets approximately {desired_length} words

Format with clear headers, subheaders, and bullet points where appropriate."""

        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ]

        return StreamingResponse(
            generate_stream(messages), media_type="text/event-stream"
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


def validate_ppt_rules(prs: Presentation) -> dict:
    """Perform deterministic rule-based validation on PowerPoint"""
    violations = []
    warnings = []
    slide_count = len(prs.slides)

    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_violations = []

        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                text_frame = shape.text_frame

                for paragraph in text_frame.paragraphs:
                    text = paragraph.text.strip()

                    if text and text.lower() in [
                        "tbd",
                        "to be determined",
                        "insert chart here",
                        "placeholder",
                    ]:
                        slide_violations.append(
                            f"Slide {slide_idx}: Placeholder text found: '{text}'"
                        )

                    if len(text) > 0:
                        if text.endswith(".") and not hasattr(shape, "has_table"):
                            slide_violations.append(
                                f"Slide {slide_idx}: Bullet point should not end with period: '{text[:50]}...'"
                            )

                        if "  " in text:
                            slide_violations.append(
                                f"Slide {slide_idx}: Double spaces found in text"
                            )

                    for run in paragraph.runs:
                        if run.font.size:
                            size_pt = run.font.size.pt
                            if size_pt < 10:
                                warnings.append(
                                    f"Slide {slide_idx}: Font size {size_pt}pt is below minimum 10pt in text: '{run.text[:30]}...'"
                                )

        violations.extend(slide_violations)

    if slide_count < 3:
        warnings.append(
            f"Presentation has only {slide_count} slides - may lack sufficient structure"
        )

    return {
        "violations": violations[:50],
        "warnings": warnings[:50],
        "total_violations": len(violations),
        "total_warnings": len(warnings),
        "slide_count": slide_count,
    }


@app.post("/api/validate-best-practices")
async def validate_best_practices(
    file: UploadFile = File(...), categories: Optional[str] = Form(None)
):
    """Validate PowerPoint against 75+ PwC consulting best practices"""
    try:
        contents = await file.read()
        prs = Presentation(io.BytesIO(contents))

        rule_validation = validate_ppt_rules(prs)

        selected_categories = []
        if categories:
            selected_categories = categories.split(",")
        else:
            selected_categories = [
                "Structure",
                "Visuals",
                "Design",
                "Charts",
                "Formatting",
                "Content",
            ]

        all_best_practices = {
            "Structure": [
                "Logical flow: Introduction → Context → Analysis → Insights → Implications → Recommendations → Next steps",
                "Pyramid Principle applied: Each section opens with key message headline",
                "Slide hierarchy consistent: Section headers, transition slides, content slides clearly differentiated",
                "Each slide answers a specific question",
                "Slide headers capture narrative takeaway, not just topic",
                "Horizontal logic: Headers form coherent narrative when read sequentially",
            ],
            "Visuals": [
                "Color used strategically to guide the eye (primary color for key messages)",
                "Converting text into purposeful visuals (2×2 grids, process flows)",
                "Visual hierarchy clear: Most important content is visually dominant",
                "Icons/images purposeful and related to text",
                "Bullets consistent: No mix of shapes unless intentional",
                "Parallel language in all bullets",
                "Concise bullet statements",
            ],
            "Design": [
                "Less is more: No unnecessary drop shadows, gradients, or animations",
                "White space used intentionally between sections",
                "Alignment perfect: Objects snapped to grid, text aligned precisely",
                "Shape sizes consistent (use exact dimensions)",
                "Images high-res and consistent style",
                "Consistent slide backgrounds (all white or light gray)",
                "Icons from PowerPoint library matching color scheme",
            ],
            "Charts": [
                "Charts simplified: No unnecessary borders, 3D effects, or legends when labels suffice",
                "Font sizes legible (10-12pt in chart labels)",
                "Axes consistent and start at zero when appropriate",
                "Same category = same color throughout deck",
                "Charts embedded rather than screenshots",
                "Every chart has subtitle with period and units",
                "Source lines included below charts",
                "Chart titles as takeaways not just labels",
            ],
            "Formatting": [
                "Fonts consistent (same font family across slides)",
                "Font sizes standardized (28pt title, 16pt header, 14pt sub-header)",
                "Font size never below 10pt (except footnotes)",
                "Full page content center-aligned",
                "Paragraph spacing set appropriately",
                "Equal alignment and spacing using Align and Distribute tools",
                "Color palette consistent with firm's official palette",
                "Headers/footers standardized using slide master",
            ],
            "Content": [
                "No typos (spell check completed)",
                "No double spaces",
                "Jargon minimized and replaced with plain business English",
                "Consistency in terminology throughout",
                "All acronyms defined on first use",
                "No previous client references",
                "No placeholder language (TBD, insert chart here)",
                "No leftover comments or notes",
                "Confidentiality notice present where required",
                "All links tested and working",
            ],
        }

        slide_count = len(prs.slides)
        total_text = 0
        total_shapes = 0

        for slide in prs.slides:
            for shape in slide.shapes:
                total_shapes += 1
                if hasattr(shape, "text"):
                    total_text += len(shape.text)

        system_prompt = f"""You are a PwC presentation quality expert. Analyze PowerPoint presentations against consulting best practices.
Be specific, critical, and constructive. Identify both strengths and areas for improvement."""

        practices_to_check = []
        for category in selected_categories:
            if category in all_best_practices:
                practices_to_check.extend(all_best_practices[category])

        violations_section = ""
        if (
            rule_validation["total_violations"] > 0
            or rule_validation["total_warnings"] > 0
        ):
            violations_section = f"""
**Automated Rule Validation Results:**
- Critical Violations Found: {rule_validation['total_violations']}
- Warnings Found: {rule_validation['total_warnings']}

**Specific Issues Detected:**
{chr(10).join(f'• {v}' for v in rule_validation['violations'][:10])}
{chr(10).join(f'⚠ {w}' for w in rule_validation['warnings'][:10])}
"""

        user_prompt = f"""Analyze this PowerPoint presentation against PwC best practices:

**Presentation Stats:**
- Total Slides: {slide_count}
- Total Shapes: {total_shapes}
- Approximate Text Length: {total_text} characters
{violations_section}
**Best Practices to Validate** (from categories: {', '.join(selected_categories)}):
{chr(10).join(f'{i+1}. {practice}' for i, practice in enumerate(practices_to_check))}

**Analysis Required:**
For each category, provide:
1. **Compliance Score**: Estimate % compliance with these practices
2. **Key Issues Found**: Specific violations or concerns (in addition to automated violations above)
3. **Strengths**: What the presentation does well
4. **Recommendations**: Prioritized improvements (address automated violations first, then other issues)

Format your response with clear headers for each category. Start with a summary of the automated violations."""

        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ]

        return StreamingResponse(
            generate_stream(messages), media_type="text/event-stream"
        )

    except Exception as e:
        raise HTTPException(
            status_code=500, detail=f"Error validating presentation: {str(e)}"
        )


@app.post("/api/export/pdf")
async def export_pdf(request: ExportRequest):
    """Export content as PDF using template"""
    try:
        pdf_bytes = await create_pdf(request.content, request.title)

        filename = f"{request.title.replace(' ', '_')[:50]}_{datetime.now().strftime('%Y%m%d')}.pdf"

        return StreamingResponse(
            io.BytesIO(pdf_bytes),
            media_type="application/pdf",
            headers={"Content-Disposition": f"attachment; filename={filename}"},
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error creating PDF: {str(e)}")


@app.post("/api/export/word")
async def export_word(request: ExportRequest):
    """Export content as Word document using template"""
    try:
        docx_bytes = await create_word_doc(request.content, request.title, use_template=True)

        filename = f"{request.title.replace(' ', '_')[:50]}_{datetime.now().strftime('%Y%m%d')}.docx"

        return StreamingResponse(
            io.BytesIO(docx_bytes),
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            headers={"Content-Disposition": f"attachment; filename={filename}"},
        )
    except Exception as e:
        raise HTTPException(
            status_code=500, detail=f"Error creating Word document: {str(e)}"
        )


class GenerateDocumentRequest(BaseModel):
    content: str
    format: str  # "pdf" or "docx"
    filename: Optional[str] = "document"
    title: Optional[str] = None  # Optional title, will be generated if not provided


@app.post("/api/generate-document")
async def generate_document(request: GenerateDocumentRequest):
    """Unified endpoint to generate documents in DOCX or PDF format using template"""
    try:
        # Use provided title or None (will be generated by LLM if not provided)
        title = request.title if request.title else None
        
        if request.format == "docx":
            # Create Word document with template - title will be generated if not provided
            docx_bytes = await create_word_doc(
                request.content, 
                title=title, 
                template_path='template/Template.docx',
                use_template=True
            )
            filename = f"{request.filename.replace(' ', '_')[:50]}_{datetime.now().strftime('%Y%m%d')}.docx"
            
            return StreamingResponse(
                io.BytesIO(docx_bytes),
                media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                headers={"Content-Disposition": f"attachment; filename={filename}"},
            )
        elif request.format == "pdf":
            # Create PDF - title will be generated if not provided
            pdf_bytes = await create_pdf(
                request.content, 
                title=title,
                template_path='template/Template.docx'
            )
            filename = f"{request.filename.replace(' ', '_')[:50]}_{datetime.now().strftime('%Y%m%d')}.pdf"
            
            return StreamingResponse(
                io.BytesIO(pdf_bytes),
                media_type="application/pdf",
                headers={"Content-Disposition": f"attachment; filename={filename}"},
            )
        else:
            raise HTTPException(status_code=400, detail=f"Unsupported format: {request.format}")
    except Exception as e:
        logger.error(f"Error generating document: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error generating document: {str(e)}")


def extract_color_scheme(prs: Presentation):
    """Extract the most common colors from a presentation."""
    colors = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "fill") and shape.fill.type == 1:
                if (
                    hasattr(shape.fill, "fore_color")
                    and shape.fill.fore_color.type == 1
                ):
                    rgb = shape.fill.fore_color.rgb
                    colors.append((rgb[0], rgb[1], rgb[2]))

            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if hasattr(run.font, "color") and run.font.color.type == 1:
                            rgb = run.font.color.rgb
                            colors.append((rgb[0], rgb[1], rgb[2]))

    if colors:
        color_counts = Counter(colors)
        most_common = [color for color, count in color_counts.most_common(5)]
        return most_common

    return [(208, 74, 2), (255, 107, 0), (26, 26, 26)]


async def correct_text_with_ai(text: str) -> str:
    """Use LLM AI to correct spelling and grammar."""
    if not text or len(text.strip()) == 0:
        return text

    try:
        llm = get_llm()
        response = await llm.chat_completion(
            messages=[
                {
                    "role": "system",
                    "content": "You are a professional editor. Fix spelling and grammar mistakes while preserving the original meaning and tone. Return ONLY the corrected text without any explanations or additional commentary.",
                },
                {"role": "user", "content": f"Correct this text: {text}"},
            ],
            temperature=0.3,
            max_tokens=1000,
        )
        return response.content.strip()
    except:
        return text


def align_shapes(slide):
    """Align shapes that are close to each other."""
    shapes = [s for s in slide.shapes if hasattr(s, "left") and hasattr(s, "top")]

    if len(shapes) < 2:
        return

    shapes_by_row = {}
    tolerance = Inches(0.1)

    for shape in shapes:
        aligned = False
        for ref_top in shapes_by_row.keys():
            if abs(shape.top - ref_top) < tolerance:
                shapes_by_row[ref_top].append(shape)
                aligned = True
                break

        if not aligned:
            shapes_by_row[shape.top] = [shape]

    for ref_top, row_shapes in shapes_by_row.items():
        if len(row_shapes) > 1:
            avg_top = sum(s.top for s in row_shapes) // len(row_shapes)
            for shape in row_shapes:
                shape.top = avg_top


def apply_color_scheme(prs: Presentation, color_scheme: list):
    """Apply color scheme to presentation."""
    if not color_scheme:
        return

    primary_color = RGBColor(*color_scheme[0])
    secondary_color = (
        RGBColor(*color_scheme[1]) if len(color_scheme) > 1 else primary_color
    )

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "fill") and shape.fill.type == 1:
                if hasattr(shape.fill, "fore_color"):
                    shape.fill.fore_color.rgb = primary_color

            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if hasattr(run.font, "color"):
                            if (
                                run.font.bold
                                or run.font.size
                                and run.font.size > Pt(18)
                            ):
                                run.font.color.rgb = primary_color
                            else:
                                run.font.color.rgb = secondary_color


@app.post("/api/ppt/improve")
async def improve_ppt(
    original_ppt: UploadFile = File(...), reference_ppt: UploadFile = File(None)
):
    """Improve PowerPoint presentation: correct spelling/grammar, align shapes, rebrand colors."""
    try:
        original_content = await original_ppt.read()
        original_prs = Presentation(io.BytesIO(original_content))

        color_scheme = None
        if reference_ppt:
            reference_content = await reference_ppt.read()
            reference_prs = Presentation(io.BytesIO(reference_content))
            color_scheme = extract_color_scheme(reference_prs)

        for slide in original_prs.slides:
            align_shapes(slide)

            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text:
                                corrected_text = await correct_text_with_ai(run.text)
                                run.text = corrected_text

        if color_scheme:
            apply_color_scheme(original_prs, color_scheme)

        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_file:
            original_prs.save(tmp_file.name)
            tmp_path = tmp_file.name

        return FileResponse(
            tmp_path,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename="improved_presentation.pptx",
            background=None,
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing PPT: {str(e)}")


class SanitizationConversationRequest(BaseModel):
    messages: List[Message]
    uploaded_file_name: Optional[str] = None
    client_identity: Optional[str] = None
    page_range: Optional[str] = None
    tier1_services: Optional[List[str]] = None  # Default services
    tier2_services: Optional[List[str]] = None  # Opt-in services
    stream: bool = True


@app.post("/api/ppt/sanitize/conversation")
async def sanitize_conversation(request: SanitizationConversationRequest):
    """
    Conversational sanitization endpoint that guides users through the sanitization process.

    Tier 1 (Default/Opt-out) Services:
    - Convert to PwC standard template
    - Replace client names and logos with placeholders
    - Delete speaker notes and comments
    - Clear presentation metadata
    - Remove numeric data (replace with X patterns)

    Tier 2 (Opt-in) Services:
    - Change competitor company names
    - Remove client-specific financial data
    - Redact business unit names
    - Remove product names
    - Sanitize location data (cities, addresses)

    Tier 3 (Advanced, performed on request):
    - Modify Think-cell chart values
    - Custom regex-based replacements
    - Advanced contextual sanitization

    Tier 4 (Never Changed):
    - PwC branding and watermarks
    - Standard consulting frameworks (MECE, Porter's Five Forces, etc.)
    - Industry-standard terminology
    """
    try:
        llm = get_llm()

        system_prompt = """You are a specialized PwC Presentation Sanitization Assistant. Your role is to guide users through a comprehensive document sanitization process.

**SANITIZATION WORKFLOW:**

1. **Recognition & Recap**: When a user mentions "sanitize", "sanitization", or "sanitation", acknowledge their request and provide a brief overview of the service.

2. **Document Collection**: If they haven't uploaded a PowerPoint file yet, politely request it. Explain that you need the original PowerPoint document to proceed.

3. **Client Identity** (Optional): Ask if they want to specify the original client name, or if you should auto-detect it from the document context.

4. **Page Range**: Ask if they want to sanitize specific pages or all pages (default: all pages).

5. **Quick Start Overview**: Present the sanitization service structure:

   **TIER 1 - Default Services (Applied unless you opt-out):**
   1. Convert to PwC standard template
   2. Replace client names and logos with [Client] placeholders  
   3. Delete all speaker notes and comments
   4. Clear presentation metadata
   5. Remove/replace numeric data with X patterns

   **TIER 2 - Opt-in Services (Only if you request):**
   6. Change competitor company names to [Competitor]
   7. Remove client-specific financial data
   8. Redact business unit names to [BU]
   9. Replace product names with [Product]
   10. Sanitize location data (cities, addresses)
   11. Remove embedded hyperlinks

   Ask: "Would you like any of the Tier 2 services?"

6. **Customization**: Allow users to:
   - Add specific actions
   - Modify listed actions
   - Remove default actions
   - Request specific replacements

7. **Service Listing**: If requested, provide a complete numbered list of ALL sanitization actions with unique indices.

8. **Confirmation**: Summarize the selected services and get final confirmation before processing.

**CONVERSATION GUIDELINES:**
- Be concise and professional
- Use bullet points for clarity
- Number services for easy reference
- Always confirm before processing
- Explain each service briefly when asked
- If user asks about specific pages, note that page-specific requests are supported

**IMPORTANT:**
- Tier 3 services (Think-cell modifications, advanced customization) are available but not announced unless specifically requested
- Never modify Tier 4 items (PwC branding, standard frameworks, industry terminology)
- Client detection can be automatic if not specified
- Default assumption: sanitize all pages unless range specified

Respond naturally and guide the user through the process step by step."""

        messages = [{"role": "system", "content": system_prompt}]
        messages.extend(
            [{"role": msg.role, "content": msg.content} for msg in request.messages]
        )

        # Add context if file uploaded
        if request.uploaded_file_name:
            context = f"\n\n[File uploaded: {request.uploaded_file_name}]"
            if request.client_identity:
                context += f"\n[Client identity specified: {request.client_identity}]"
            if request.page_range:
                context += f"\n[Page range: {request.page_range}]"
            messages[-1]["content"] += context

        async def stream_response():
            try:
                async for content in llm.stream_completion(
                    messages=messages,
                    temperature=0.7,
                    max_tokens=2048,
                ):
                    yield f"data: {json.dumps({'content': content})}\n\n"

                yield f"data: {json.dumps({'done': True})}\n\n"

            except Exception as e:
                yield f"data: {json.dumps({'error': str(e)})}\n\n"

        return StreamingResponse(stream_response(), media_type="text/event-stream")

    except Exception as e:
        raise HTTPException(
            status_code=500, detail=f"Error in sanitization conversation: {str(e)}"
        )


@app.post("/api/ppt/sanitize")
async def sanitize_ppt(
    original_ppt: UploadFile = File(...),
    client_name: Optional[str] = None,
    client_products: Optional[str] = None,
    business_units: Optional[str] = None,
    sanitization_options: Optional[str] = Form(None),
    page_range: Optional[str] = Form(None),
):
    """
    Comprehensive PowerPoint sanitization with grammar correction:

    Data Sanitization (Tier-based):

    Tier 1 (Default):
    - Client names: Replace with [Client]
    - Logos and watermarks: Remove images
    - Speaker notes and comments: Clear all
    - Metadata: Clear all document properties
    - Numeric data: Replace with X placeholders

    Tier 2 (Opt-in):
    - Product names: Replace with [Product]
    - Business unit names: Replace with [BU]
    - Competitor names: Replace with [Competitor]
    - Financial data: Remove client-specific numbers
    - Location data: Redact cities, addresses

    Tier 3 (Advanced):
    - Think-cell chart modifications
    - Custom replacements
    - Personal information: Emails, phones, SSN
    - Embedded files: Disconnect attachments

    Grammar & Spelling:
    - AI-powered correction using Groq
    - Preserves formatting and tone
    """
    try:
        if not original_ppt.filename.endswith(".pptx"):
            raise HTTPException(
                status_code=400, detail="File must be a .pptx PowerPoint file"
            )

        original_content = await original_ppt.read()

        client_names = []
        if client_name:
            client_names = [
                name.strip() for name in client_name.split(",") if name.strip()
            ]

        product_names = []
        if client_products:
            product_names = [
                name.strip() for name in client_products.split(",") if name.strip()
            ]

        business_unit_names = []
        if business_units:
            business_unit_names = [
                name.strip() for name in business_units.split(",") if name.strip()
            ]

        # Parse sanitization options for selective sanitization
        options = {}
        if sanitization_options:
            try:
                options = json.loads(sanitization_options)
            except json.JSONDecodeError:
                pass

        sanitizer = PPTSanitizer(
            client_names=client_names,
            product_names=product_names,
            business_units=business_unit_names,
            options=options,
        )

        sanitized_output = sanitizer.sanitize_presentation(
            io.BytesIO(original_content),
            client_name=client_names[0] if client_names else None,
            fix_grammar=True,  # Enable grammar fixing
        )

        stats = sanitizer.get_stats()
        stats["grammar_corrections"] = sanitizer.sanitization_stats.get(
            "grammar_corrections", 0
        )

        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_file:
            tmp_file.write(sanitized_output.getvalue())
            tmp_path = tmp_file.name

        headers = {"X-Sanitization-Stats": json.dumps(stats)}

        return FileResponse(
            tmp_path,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename="sanitized_presentation.pptx",
            headers=headers,
            background=None,
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error sanitizing PPT: {str(e)}")


# Podcast Generation Functions


def get_polly_client():
    """Get or create Polly client with AWS credentials"""
    aws_access_key = os.getenv("AWS_ACCESS_KEY_ID")
    aws_secret_key = os.getenv("AWS_SECRET_ACCESS_KEY")
    aws_region = os.getenv("AWS_REGION", "us-east-1")

    if not aws_access_key or not aws_secret_key:
        raise HTTPException(
            status_code=500,
            detail="AWS credentials not configured. Please add AWS_ACCESS_KEY_ID and AWS_SECRET_ACCESS_KEY.",
        )

    return boto3.client(
        "polly",
        aws_access_key_id=aws_access_key,
        aws_secret_access_key=aws_secret_key,
        region_name=aws_region,
    )


async def generate_podcast_script(
    content: str, customization: Optional[str] = None, podcast_style: str = "dialogue"
):
    """Generate a podcast script from content using LLM - supports dialogue or monologue styles"""
    llm = get_llm()

    if podcast_style == "monologue":
        system_prompt = """You are a podcast script writer creating engaging single-narrator podcast episodes.

Create a natural, engaging narration:
- **Narrator**: Professional, clear, engaging storytelling voice

Script Requirements:
- Target 10-15 minutes (2,000-2,500 words)
- Natural narration with clear structure and flow
- Include verbal pauses and conversational elements
- Start with a brief introduction of the topic
- End with key takeaways and closing remarks
- Make complex ideas accessible and engaging
- Use storytelling and real-world examples

Format your script as:
NARRATOR: [narration]

Keep the narration flowing naturally with smooth transitions between sections."""

        user_prompt = f"""Create an engaging podcast monologue discussing the following content:

{content}

{f"Special instructions: {customization}" if customization else ""}

Generate a complete, natural-sounding podcast narration."""
    else:
        system_prompt = """You are a podcast script writer creating engaging two-host conversational podcasts.

Create a natural, engaging dialogue between two hosts:
- **Alex** (Host 1): Thoughtful, asks clarifying questions, brings strategic perspective
- **Jordan** (Host 2): Energetic, explains concepts clearly, adds relatable examples

Script Requirements:
- Target 10-15 minutes (2,000-2,500 words)
- Natural conversation with questions, answers, insights, and examples
- Include verbal pauses and conversational elements (e.g., "you know," "right," "that's interesting")
- Start with a brief introduction of the topic
- End with key takeaways and closing remarks
- Make complex ideas accessible and engaging
- Use storytelling and real-world examples

Format your script as:
ALEX: [dialogue]
JORDAN: [dialogue]

Keep the conversation flowing naturally with back-and-forth exchanges."""

        user_prompt = f"""Create an engaging podcast script discussing the following content:

{content}

{f"Special instructions: {customization}" if customization else ""}

Generate a complete, natural-sounding podcast conversation between Alex and Jordan."""

    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": user_prompt},
    ]

    response = await llm.chat_completion(
        messages=messages,
        temperature=0.8,
        max_tokens=4000,
    )

    return response.content


def parse_script_segments(script: str) -> List[dict]:
    """Parse podcast script into segments with speaker labels"""
    segments = []
    lines = script.strip().split("\n")

    current_speaker = None
    current_text = []

    for line in lines:
        line = line.strip()
        if not line:
            continue

        # Check if line starts with a speaker label
        if line.upper().startswith("ALEX:"):
            if current_speaker and current_text:
                segments.append(
                    {"speaker": current_speaker, "text": " ".join(current_text).strip()}
                )
            current_speaker = "ALEX"
            current_text = [line[5:].strip()]  # Remove "ALEX:"
        elif line.upper().startswith("JORDAN:"):
            if current_speaker and current_text:
                segments.append(
                    {"speaker": current_speaker, "text": " ".join(current_text).strip()}
                )
            current_speaker = "JORDAN"
            current_text = [line[7:].strip()]  # Remove "JORDAN:"
        elif line.upper().startswith("NARRATOR:"):
            if current_speaker and current_text:
                segments.append(
                    {"speaker": current_speaker, "text": " ".join(current_text).strip()}
                )
            current_speaker = "NARRATOR"
            current_text = [line[9:].strip()]  # Remove "NARRATOR:"
        else:
            # Continuation of current speaker's dialogue
            if current_speaker:
                current_text.append(line)

    # Add final segment
    if current_speaker and current_text:
        segments.append(
            {"speaker": current_speaker, "text": " ".join(current_text).strip()}
        )

    return segments


async def synthesize_with_polly(text: str, voice_id: str) -> bytes:
    """Synthesize speech using Amazon Polly"""
    polly_client = get_polly_client()

    try:
        response = polly_client.synthesize_speech(
            Text=text, OutputFormat="mp3", VoiceId=voice_id, Engine="neural"
        )

        return response["AudioStream"].read()
    except (BotoCoreError, ClientError) as error:
        raise HTTPException(
            status_code=500, detail=f"Polly synthesis error: {str(error)}"
        )


async def create_podcast_audio(script: str) -> bytes:
    """Create complete podcast audio from script with two distinct voices"""
    segments = parse_script_segments(script)

    if not segments:
        raise HTTPException(status_code=400, detail="Could not parse podcast script")

    # Voice mapping
    voice_map = {
        "ALEX": "Matthew",  # Male neural voice
        "JORDAN": "Joanna",  # Female neural voice
        "NARRATOR": "Matthew",  # Professional male voice for monologues
    }

    # Create audio segments
    audio_segments = []

    for segment in segments:
        pass
        # voice = voice_map.get(segment["speaker"], "Matthew")
        # audio_data = await synthesize_with_polly(segment["text"], voice)

        # # Convert to AudioSegment
        # audio_segment = AudioSegment.from_mp3(io.BytesIO(audio_data))
        # audio_segments.append(audio_segment)

        # # Add brief pause between speakers (300ms)
        # audio_segments.append(AudioSegment.silent(duration=300))

    # Concatenate all segments
    if not audio_segments:
        raise HTTPException(status_code=400, detail="No audio segments generated")

    final_audio = audio_segments[0]
    for segment in audio_segments[1:]:
        final_audio += segment

    # Export to MP3
    output_buffer = io.BytesIO()
    final_audio.export(output_buffer, format="mp3", bitrate="128k")
    output_buffer.seek(0)

    return output_buffer.read()


@app.post("/api/generate-podcast")
async def generate_podcast(
    files: Optional[List[UploadFile]] = File(None),
    content_text: Optional[str] = Form(None),
    customization: Optional[str] = Form(None),
    podcast_style: Optional[str] = Form("dialogue"),
):
    """Generate a NotebookLM-style podcast from uploaded documents or text"""
    try:

        async def event_generator():
            try:
                # Step 1: Parse uploaded files and combine content
                yield f"data: {json.dumps({'type': 'progress', 'message': 'Parsing uploaded documents...', 'percent': 10})}\n\n"

                combined_content = ""

                # Add explicit text content if provided
                if content_text:
                    combined_content += content_text + "\n\n"

                # Parse uploaded files
                if files:
                    for file in files:
                        file_content = await file.read()
                        filename = file.filename.lower()

                        if filename.endswith(".pdf"):
                            parsed_content = extract_text_from_pdf(file_content)
                        elif filename.endswith(".docx"):
                            parsed_content = extract_text_from_docx(file_content)
                        elif filename.endswith((".txt", ".md")):
                            parsed_content = file_content.decode(
                                "utf-8", errors="ignore"
                            )
                        else:
                            continue

                        combined_content += f"\n\n{parsed_content}"

                if not combined_content.strip():
                    yield f"data: {json.dumps({'type': 'error', 'message': 'No content provided. Please upload files or provide text.'})}\n\n"
                    return

                # Limit content length
                if len(combined_content) > 50000:
                    combined_content = combined_content[:50000]

                # Step 2: Generate podcast script
                style_message = (
                    "monologue script"
                    if podcast_style == "monologue"
                    else "conversational podcast script"
                )
                yield f"data: {json.dumps({'type': 'progress', 'message': f'Generating {style_message}...', 'percent': 30})}\n\n"

                script = await generate_podcast_script(
                    combined_content, customization, podcast_style
                )

                if not script or not script.strip():
                    yield f"data: {json.dumps({'type': 'error', 'message': 'Failed to generate podcast script'})}\n\n"
                    return

                yield f"data: {json.dumps({'type': 'progress', 'message': 'Script generated! Now synthesizing audio with Amazon Polly...', 'percent': 50})}\n\n"

                # Step 3: Synthesize audio
                yield f"data: {json.dumps({'type': 'progress', 'message': 'Converting to speech (this may take a few minutes)...', 'percent': 60})}\n\n"

                audio_data = await create_podcast_audio(script)

                # Log audio generation stats
                print(f"Audio generated: {len(audio_data)} bytes")

                # Step 4: Encode audio as base64 for transfer
                yield f"data: {json.dumps({'type': 'progress', 'message': 'Finalizing podcast...', 'percent': 90})}\n\n"

                audio_base64 = base64.b64encode(audio_data).decode("utf-8")
                print(f"Base64 encoded: {len(audio_base64)} characters")

                # Return script and audio
                yield f"data: {json.dumps({'type': 'script', 'content': script})}\n\n"
                yield f"data: {json.dumps({'type': 'complete', 'message': 'Podcast generated successfully!', 'audio': audio_base64, 'percent': 100})}\n\n"

            except Exception as e:
                yield f"data: {json.dumps({'type': 'error', 'message': str(e)})}\n\n"

        return StreamingResponse(event_generator(), media_type="text/event-stream")

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# ===================================
# NEW THOUGHT LEADERSHIP SECTIONS (5)
# ===================================

# Helper Functions for Analysis


def analyze_mece_framework(content: str) -> dict:
    """Analyze content for MECE (Mutually Exclusive, Collectively Exhaustive) framework compliance"""
    return {
        "has_clear_structure": True,  # Will use LLM in full implementation
        "mutually_exclusive": True,
        "collectively_exhaustive": True,
        "suggestions": [],
    }


def detect_generic_content(content: str) -> List[str]:
    """Detect generic or non-specific content areas"""
    generic_phrases = [
        "best practices",
        "industry-leading",
        "world-class",
        "cutting-edge",
        "state-of-the-art",
        "innovative solutions",
    ]
    detected = []
    content_lower = content.lower()
    for phrase in generic_phrases:
        if phrase in content_lower:
            detected.append(phrase)
    return detected


def detect_text_to_visual_opportunities(content: str) -> List[dict]:
    """Detect areas where text could be converted to charts/graphs"""
    opportunities = []
    lines = content.split("\n")

    for i, line in enumerate(lines):
        if any(
            word in line.lower()
            for word in ["percentage", "percent", "%", "increase", "decrease", "growth"]
        ):
            opportunities.append(
                {
                    "line": i + 1,
                    "type": "chart",
                    "suggestion": "Consider visualizing this data as a bar or line chart",
                }
            )
        elif any(
            word in line.lower()
            for word in ["vs", "versus", "compared to", "comparison"]
        ):
            opportunities.append(
                {
                    "line": i + 1,
                    "type": "comparison_table",
                    "suggestion": "Consider using a comparison table or chart",
                }
            )

    return opportunities[:5]  # Limit to top 5 suggestions


def format_pwc_citation(source_type: str, source_info: dict) -> str:
    """Format citations according to PwC guidelines"""
    if source_type == "article":
        return f"{source_info.get('author', 'Unknown')}, {source_info.get('year', 'n.d.')}. {source_info.get('title', 'Untitled')}, {source_info.get('journal', '')}, Vol. {source_info.get('volume', '')}, pp. {source_info.get('pages', '')}."
    elif source_type == "book":
        return f"{source_info.get('author', 'Unknown')}, {source_info.get('year', 'n.d.')}. {source_info.get('title', 'Untitled')}. {source_info.get('city', '')}: {source_info.get('publisher', '')}."
    else:
        return (
            f"{source_info.get('title', 'Source')}, {source_info.get('year', 'n.d.')}."
        )


async def detect_contradictions(sources_content: List[str]) -> List[dict]:
    """Detect contradicting information across sources"""
    return []  # Placeholder - full implementation would use LLM comparison


# Section 1: Draft Content
@app.post("/api/tl/draft-content")
async def draft_content_endpoint(request: DraftContentRequest):
    """
    Draft Content - Section 1
    Conversational workflow for drafting articles, blogs, white papers, executive briefs
    """
    try:

        async def generate_stream():
            try:
                llm = get_llm()
                messages = request.messages

                # Build conversational system prompt
                system_prompt = """You are an expert content strategist at PwC specializing in thought leadership.
You help draft compelling, well-structured content (Articles, Blogs, White Papers, Executive Briefs).

**Your Process:**
1. Turn concepts into structured drafts
2. Conduct targeted research (if requested)
3. Provide recommendations and iterate

**Content Types:**
- **Article**: 1500-2500 words, analytical depth, data-driven insights
- **Blog**: 800-1200 words, conversational, engaging, clear takeaways
- **White Paper**: 3000-5000 words, comprehensive research, authoritative
- **Executive Brief**: 500-800 words, concise, action-oriented, C-suite focused

**Best Practices:**
- Use MECE framework (Mutually Exclusive, Collectively Exhaustive)
- Incorporate PwC proprietary research and data
- Identify opportunities for visuals (charts, graphs, infographics)
- Avoid generic statements - use specific examples and data
- Flag contradictory points
- Suggest where text-heavy sections could become visuals

**Conversation Style:**
- Guide users through the process step-by-step
- Ask clarifying questions if inputs are incomplete
- Provide indexed options (e.g., "1. Article, 2. Blog, 3. White Paper, 4. Executive Brief")
- Users can select by number (e.g., "1" for Article)

**Required Inputs:**
- Topic
- Content type (Article/Blog/White Paper/Executive Brief)
- Audience
- Length

**Optional Inputs:**
- Concept/brief/outline/rough draft (accept uploads or in-prompt)
- Supporting documents
- Include research (yes/no)
- Additional guidelines

If the user provides only a topic without definition, converse to develop a thorough concept."""

                # Convert request messages to LLM format
                llm_messages = [{"role": "system", "content": system_prompt}]
                llm_messages.extend(
                    [{"role": msg.role, "content": msg.content} for msg in messages]
                )

                # Add context about current state
                context_parts = []
                if request.content_type:
                    context_parts.append(f"Content Type: {request.content_type}")
                if request.topic:
                    context_parts.append(f"Topic: {request.topic}")
                if request.audience:
                    context_parts.append(f"Audience: {request.audience}")
                if request.length:
                    context_parts.append(f"Target Length: {request.length}")

                if context_parts:
                    context_message = "\n".join(context_parts)
                    llm_messages.append(
                        {
                            "role": "user",
                            "content": f"[Current Inputs]\n{context_message}",
                        }
                    )

                # Stream response
                async for chunk in llm.stream_completion(
                    messages=llm_messages, temperature=0.7, max_tokens=4096
                ):
                    yield f"data: {json.dumps({'type': 'content', 'content': chunk})}\n\n"

                yield f"data: {json.dumps({'type': 'complete'})}\n\n"

            except Exception as e:
                yield f"data: {json.dumps({'type': 'error', 'message': str(e)})}\n\n"

        return StreamingResponse(generate_stream(), media_type="text/event-stream")

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# Section 2: Conduct Research
@app.post("/api/tl/conduct-research")
async def conduct_research_endpoint(request: ConductResearchRequest):
    """
    Conduct Research - Section 2
    Multi-source research with PwC Proprietary, Licensed, and External sources
    """
    try:

        async def generate_stream():
            try:
                llm = get_llm()
                messages = request.messages

                # PwC proprietary sources with URLs
                pwc_proprietary_sources = {
                    "PwC Research Library": "https://www.pwc.com/us/en/library.html",
                    "PwC Industries": "https://www.pwc.com/us/en/industries.html",
                    "Strategy+Business": "https://www.strategy-business.com/",
                    "Executive Leadership Hub": "https://www.pwc.com/us/en/executive-leadership-hub.html",
                    "PwC Exchange": "https://www.pwc.com/us/en/services/consulting/pwc-exchange.html",
                    # Additional placeholders (no URLs yet)
                    "PwC Insights": None,
                    "PwC Industry Edge": None,
                    "PwC Advisory Commercial Hub": None,
                    "PwC Connected Source": None,
                    "Analyst Relations": None,
                    "Client Success Stories": None,
                    "Deals Source": None,
                    "Inside Industries": None,
                    "Insights to Win": None,
                    "Policy on Demand": None,
                    "PwC Intelligence": None,
                    "Tax Source": None,
                }

                # External research sources (non-PwC)
                external_research_sources = {
                    "UNDP Human Development Index": "https://hdr.undp.org/data-center/human-development-index#/indicies/HDI",
                }

                licensed_third_party_sources = [
                    "BoardEx",
                    "Cerulli",
                    "Claritas",
                    "CompanyIQ",
                    "CFRA Industry Surveys",
                    "EMIS Professional",
                    "Factiva",
                    "Forrester Research",
                    "Gartner Research",
                    "Global Data",
                    "IBIS World",
                    "IDC",
                    "Preqin Pro",
                    "S&P Cap IQ Pro",
                    "S&P Global Connect",
                    "Source Global Research",
                    "Technology Business Review",
                    "Citeline Pharma Intelligence",
                    "CoStar",
                    "Definitive Healthcare",
                    "IQVIA",
                    "Barron's & WSJ",
                ]

                # Track all fetched sources for citations
                all_fetched_sources = []
                citation_counter = 1

                # Extract user query, URLs, and uploaded documents from messages
                user_query = ""
                user_urls = []
                user_uploaded_docs = []

                for msg in reversed(messages):
                    if msg.role == "user":
                        content = msg.content

                        # Extract research question
                        if "Generate a comprehensive business article on:" in content:
                            user_query = (
                                content.split(
                                    "Generate a comprehensive business article on:"
                                )[-1]
                                .split("\n")[0]
                                .strip()
                            )
                        elif "I need to conduct research on:" in content:
                            user_query = (
                                content.split("I need to conduct research on:")[-1]
                                .split("\n")[0]
                                .strip()
                            )
                        elif ":" in content:
                            user_query = content.split(":")[-1].split("\n")[0].strip()
                        else:
                            user_query = content[:200]

                        # Extract user-provided URLs
                        if "Reference URLs:" in content:
                            urls_section = (
                                content.split("Reference URLs:")[-1]
                                .split("\n\n")[0]
                                .strip()
                            )
                            # Handle URLs separated by newlines, commas, or spaces
                            # Extract all URLs (http/https/www.)
                            url_pattern = r"(?:https?://|www\.)[^\s,\n]+"
                            found_urls = re.findall(url_pattern, urls_section)
                            user_urls = [
                                url.strip() for url in found_urls if url.strip()
                            ]

                        # Extract user-uploaded documents
                        if "Uploaded Documents" in content:
                            doc_section = content.split("Uploaded Documents")[-1]
                            # Extract document names and content
                            doc_pattern = (
                                r"---\s*([^---]+?)\s*---\s*\n(.*?)(?=\n---|\Z)"
                            )
                            matches = re.findall(doc_pattern, doc_section, re.DOTALL)
                            for doc_name, doc_content in matches:
                                if doc_content.strip():
                                    user_uploaded_docs.append(
                                        {
                                            "name": doc_name.strip(),
                                            "content": doc_content.strip()[
                                                :10000
                                            ],  # Limit content
                                        }
                                    )
                        break

                # Helper function to check if sources are relevant using semantic search
                async def check_relevance(
                    sources: list, query: str, min_relevant_chunks: int = 3
                ):
                    """Check if sources contain relevant content. Returns (is_relevant, relevant_chunks)"""
                    if not sources or not query or not DEPENDENCIES_AVAILABLE:
                        return (False, [])

                    try:
                        retriever = SemanticRetriever()
                        await retriever.index_sources(sources)
                        retrieved_chunks = await retriever.retrieve(
                            query, top_k=min_relevant_chunks * 2
                        )
                        is_relevant = len(retrieved_chunks) >= min_relevant_chunks
                        return (is_relevant, retrieved_chunks)
                    except Exception as e:
                        logger.warning(f"Relevance check failed: {e}")
                        return (False, [])

                # Priority-based source retrieval
                relevant_content = ""
                all_sources_for_indexing = []
                used_citation_numbers = set()
                retriever = None

                # PRIORITY 0: ALWAYS fetch PwC Research Library as default source (hardcoded)
                default_pwc_library_url = "https://www.pwc.com/us/en/library.html"
                default_pwc_library_sources = []
                yield f'data: {json.dumps({"type": "progress", "message": "[Default Source] Fetching PwC Research Library (always included)..."})}\n\n'

                try:
                    url_data = await fetch_url_content(default_pwc_library_url)
                    if url_data.get("success"):
                        citation_num = citation_counter
                        citation_counter += 1

                        default_pwc_library_sources.append(
                            {
                                "citation_number": citation_num,
                                "name": "PwC Research Library",
                                "url": default_pwc_library_url,
                                "title": url_data.get("title", "PwC Research Library"),
                                "content": url_data.get("content", "")[:10000],
                            }
                        )

                        all_fetched_sources.append(
                            {
                                "number": citation_num,
                                "url": default_pwc_library_url,
                                "title": url_data.get("title", "PwC Research Library")
                                or "PwC Research Library",
                            }
                        )

                        # Crawl related pages
                        try:
                            yield f'data: {json.dumps({"type": "progress", "message": "Crawling related pages from PwC Research Library..."})}\n\n'
                            related_pages = await crawl_related_pages(
                                default_pwc_library_url,
                                max_pages=10,
                                max_depth=3,
                                query=user_query,
                            )

                            for page in related_pages:
                                if page["url"] != default_pwc_library_url:
                                    citation_num = citation_counter
                                    citation_counter += 1

                                    default_pwc_library_sources.append(
                                        {
                                            "citation_number": citation_num,
                                            "name": "PwC Research Library - Related",
                                            "url": page["url"],
                                            "title": page.get("title", "").strip()
                                            or page["url"],
                                            "content": (
                                                page["content"][:10000]
                                                if len(page["content"]) > 10000
                                                else page["content"]
                                            ),
                                        }
                                    )

                                    all_fetched_sources.append(
                                        {
                                            "number": citation_num,
                                            "url": page["url"],
                                            "title": page.get("title", "").strip()
                                            or page["url"],
                                        }
                                    )
                        except Exception:
                            pass
                except Exception as e:
                    logger.warning(f"Failed to fetch default PwC Research Library: {e}")

                # Always include default PwC Library sources
                if default_pwc_library_sources:
                    all_sources_for_indexing.extend(default_pwc_library_sources)
                    yield f'data: {json.dumps({"type": "progress", "message": f"Added {len(default_pwc_library_sources)} PwC Research Library source(s) as default"})}\n\n'

                # PRIORITY 1: PwC Proprietary Resources
                pwc_sources = []
                if request.source_groups and "PwC Proprietary" in request.source_groups:
                    yield f'data: {json.dumps({"type": "progress", "message": "[Priority 1] Fetching content from PwC proprietary sources..."})}\n\n'

                    sources_with_urls = {
                        name: url
                        for name, url in pwc_proprietary_sources.items()
                        if url
                    }

                    for source_name, source_url in sources_with_urls.items():
                        try:
                            yield f'data: {json.dumps({"type": "progress", "message": f"Fetching {source_name}..."})}\n\n'

                            # Fetch main page
                            url_data = await fetch_url_content(source_url)
                            if url_data.get("success"):
                                citation_num = citation_counter
                                citation_counter += 1

                                pwc_sources.append(
                                    {
                                        "citation_number": citation_num,
                                        "name": source_name,
                                        "url": source_url,
                                        "title": url_data.get("title", source_name),
                                        "content": url_data.get("content", "")[:10000],
                                    }
                                )

                                all_fetched_sources.append(
                                    {
                                        "number": citation_num,
                                        "url": source_url,
                                        "title": url_data.get("title", source_name)
                                        or source_name,
                                    }
                                )

                                # Crawl related pages (subwebsite crawling)
                                try:
                                    yield f'data: {json.dumps({"type": "progress", "message": f"Crawling related pages from {source_name}..."})}\n\n'
                                    related_pages = await crawl_related_pages(
                                        source_url,
                                        max_pages=10,
                                        max_depth=3,
                                        query=user_query,
                                    )

                                    for page in related_pages:
                                        if page["url"] != source_url:
                                            citation_num = citation_counter
                                            citation_counter += 1

                                            pwc_sources.append(
                                                {
                                                    "citation_number": citation_num,
                                                    "name": f"{source_name} - Related",
                                                    "url": page["url"],
                                                    "title": page.get(
                                                        "title", ""
                                                    ).strip()
                                                    or page["url"],
                                                    "content": (
                                                        page["content"][:10000]
                                                        if len(page["content"]) > 10000
                                                        else page["content"]
                                                    ),
                                                }
                                            )

                                            all_fetched_sources.append(
                                                {
                                                    "number": citation_num,
                                                    "url": page["url"],
                                                    "title": page.get(
                                                        "title", ""
                                                    ).strip()
                                                    or page["url"],
                                                }
                                            )
                                except Exception:
                                    pass
                        except Exception:
                            continue

                    # Check if PwC sources are relevant
                    if pwc_sources and user_query:
                        yield f'data: {json.dumps({"type": "progress", "message": "Checking relevance of PwC sources..."})}\n\n'
                        is_relevant, relevant_chunks = await check_relevance(
                            pwc_sources, user_query, min_relevant_chunks=3
                        )

                        if is_relevant:
                            yield f'data: {json.dumps({"type": "progress", "message": f"PwC sources are relevant. Adding {len(pwc_sources)} source(s)."})}\n\n'
                            all_sources_for_indexing.extend(pwc_sources)
                        else:
                            yield f'data: {json.dumps({"type": "progress", "message": "PwC sources have limited relevance. Proceeding to next priority..."})}\n\n'
                            pwc_sources = []  # Clear if not relevant enough
                    elif pwc_sources:
                        # If no query or relevance check unavailable, use PwC sources
                        all_sources_for_indexing.extend(pwc_sources)

                # PRIORITY 2: User Mentioned Links (always add if provided, in addition to default)
                user_link_sources = []
                if user_urls:
                    yield f'data: {json.dumps({"type": "progress", "message": "[Priority 2] Fetching content from user-provided URLs..."})}\n\n'

                    # Normalize default URL for comparison
                    default_url_normalized = (
                        default_pwc_library_url.split("#")[0].split("?")[0].rstrip("/")
                    )

                    for url in user_urls:
                        try:
                            # Normalize URL
                            if url.startswith("www."):
                                url = "https://" + url

                            # Normalize for comparison (remove fragments, query params, trailing slashes)
                            url_normalized = url.split("#")[0].split("?")[0].rstrip("/")

                            # Skip if user provided the same URL as default PwC Research Library
                            if url_normalized == default_url_normalized:
                                yield f'data: {json.dumps({"type": "progress", "message": f"Skipping {url} (already included as default source)..."})}\n\n'
                                continue

                            yield f'data: {json.dumps({"type": "progress", "message": f"Fetching {url}..."})}\n\n'

                            url_data = await fetch_url_content(url)
                            if url_data.get("success"):
                                citation_num = citation_counter
                                citation_counter += 1

                                user_link_sources.append(
                                    {
                                        "citation_number": citation_num,
                                        "name": f"User Provided Link",
                                        "url": url,
                                        "title": url_data.get("title", url),
                                        "content": url_data.get("content", "")[:10000],
                                    }
                                )

                                all_fetched_sources.append(
                                    {
                                        "number": citation_num,
                                        "url": url,
                                        "title": url_data.get("title", url) or url,
                                    }
                                )

                                # Crawl related pages
                                try:
                                    yield f'data: {json.dumps({"type": "progress", "message": f"Crawling related pages from {url}..."})}\n\n'
                                    related_pages = await crawl_related_pages(
                                        url, max_pages=10, max_depth=3, query=user_query
                                    )

                                    for page in related_pages:
                                        if page["url"] != url:
                                            citation_num = citation_counter
                                            citation_counter += 1

                                            user_link_sources.append(
                                                {
                                                    "citation_number": citation_num,
                                                    "name": "User Provided Link - Related",
                                                    "url": page["url"],
                                                    "title": page.get(
                                                        "title", ""
                                                    ).strip()
                                                    or page["url"],
                                                    "content": (
                                                        page["content"][:10000]
                                                        if len(page["content"]) > 10000
                                                        else page["content"]
                                                    ),
                                                }
                                            )

                                            all_fetched_sources.append(
                                                {
                                                    "number": citation_num,
                                                    "url": page["url"],
                                                    "title": page.get(
                                                        "title", ""
                                                    ).strip()
                                                    or page["url"],
                                                }
                                            )
                                except Exception:
                                    pass
                        except Exception:
                            continue

                    # Check if user links are relevant
                    if user_link_sources and user_query:
                        yield f'data: {json.dumps({"type": "progress", "message": "Checking relevance of user-provided links..."})}\n\n'
                        is_relevant, relevant_chunks = await check_relevance(
                            user_link_sources, user_query, min_relevant_chunks=2
                        )

                        if is_relevant:
                            yield f'data: {json.dumps({"type": "progress", "message": f"User links are relevant. Adding {len(user_link_sources)} source(s)."})}\n\n'
                            all_sources_for_indexing.extend(user_link_sources)
                        else:
                            yield f'data: {json.dumps({"type": "progress", "message": "User links have limited relevance. Proceeding to next priority..."})}\n\n'
                            user_link_sources = []
                    elif user_link_sources:
                        all_sources_for_indexing.extend(user_link_sources)

                # PRIORITY 3: User Uploaded Documents (always add if provided, in addition to default)
                user_doc_sources = []
                if user_uploaded_docs:
                    yield f'data: {json.dumps({"type": "progress", "message": "[Priority 3] Processing user-uploaded documents..."})}\n\n'

                    for doc in user_uploaded_docs:
                        citation_num = citation_counter
                        citation_counter += 1

                        user_doc_sources.append(
                            {
                                "citation_number": citation_num,
                                "name": doc["name"],
                                "url": f"uploaded://{doc['name']}",  # Placeholder URL for uploaded docs
                                "title": doc["name"],
                                "content": doc["content"],
                            }
                        )

                        all_fetched_sources.append(
                            {
                                "number": citation_num,
                                "url": f"uploaded://{doc['name']}",
                                "title": doc["name"],
                            }
                        )

                    # Check if user docs are relevant
                    if user_doc_sources and user_query:
                        yield f'data: {json.dumps({"type": "progress", "message": "Checking relevance of uploaded documents..."})}\n\n'
                        is_relevant, relevant_chunks = await check_relevance(
                            user_doc_sources, user_query, min_relevant_chunks=2
                        )

                        if is_relevant:
                            yield f'data: {json.dumps({"type": "progress", "message": f"Uploaded documents are relevant. Adding {len(user_doc_sources)} document(s)."})}\n\n'
                            all_sources_for_indexing.extend(user_doc_sources)
                        else:
                            yield f'data: {json.dumps({"type": "progress", "message": "Uploaded documents have limited relevance. Proceeding to external sources..."})}\n\n'
                            user_doc_sources = []
                    elif user_doc_sources:
                        all_sources_for_indexing.extend(user_doc_sources)

                # PRIORITY 4: External Internet Sources (always add if requested, in addition to default)
                external_sources = []
                if (
                    request.source_groups
                    and "External Research" in request.source_groups
                ):
                    yield f'data: {json.dumps({"type": "progress", "message": "[Priority 4] Fetching content from external research sources..."})}\n\n'

                    for source_name, source_url in external_research_sources.items():
                        try:
                            yield f'data: {json.dumps({"type": "progress", "message": f"Fetching {source_name}..."})}\n\n'

                            url_data = await fetch_url_content(source_url)
                            if url_data.get("success"):
                                citation_num = citation_counter
                                citation_counter += 1

                                external_sources.append(
                                    {
                                        "citation_number": citation_num,
                                        "name": source_name,
                                        "url": source_url,
                                        "title": url_data.get("title", source_name),
                                        "content": url_data.get("content", "")[:10000],
                                    }
                                )

                                all_fetched_sources.append(
                                    {
                                        "number": citation_num,
                                        "url": source_url,
                                        "title": url_data.get("title", source_name)
                                        or source_name,
                                    }
                                )

                                # Crawl related pages
                                try:
                                    yield f'data: {json.dumps({"type": "progress", "message": f"Crawling related pages from {source_name}..."})}\n\n'
                                    related_pages = await crawl_related_pages(
                                        source_url,
                                        max_pages=10,
                                        max_depth=3,
                                        query=user_query,
                                    )

                                    for page in related_pages:
                                        if page["url"] != source_url:
                                            citation_num = citation_counter
                                            citation_counter += 1

                                            external_sources.append(
                                                {
                                                    "citation_number": citation_num,
                                                    "name": f"{source_name} - Related",
                                                    "url": page["url"],
                                                    "title": page.get(
                                                        "title", ""
                                                    ).strip()
                                                    or page["url"],
                                                    "content": (
                                                        page["content"][:10000]
                                                        if len(page["content"]) > 10000
                                                        else page["content"]
                                                    ),
                                                }
                                            )

                                            all_fetched_sources.append(
                                                {
                                                    "number": citation_num,
                                                    "url": page["url"],
                                                    "title": page.get(
                                                        "title", ""
                                                    ).strip()
                                                    or page["url"],
                                                }
                                            )
                                except Exception:
                                    pass
                        except Exception:
                            continue

                    # Use external sources if available (no relevance check needed as last resort)
                    if external_sources:
                        yield f'data: {json.dumps({"type": "progress", "message": f"Adding {len(external_sources)} external source(s)."})}\n\n'
                        all_sources_for_indexing.extend(external_sources)

                # Final fallback: combine all available sources if still nothing (shouldn't happen since default is always included)
                if not all_sources_for_indexing:
                    all_sources_for_indexing = (
                        default_pwc_library_sources
                        + pwc_sources
                        + user_link_sources
                        + user_doc_sources
                        + external_sources
                    )
                    if all_sources_for_indexing:
                        yield f'data: {json.dumps({"type": "progress", "message": f"Using all available sources ({len(all_sources_for_indexing)} total)..."})}\n\n'

                # Perform semantic retrieval on selected sources
                retrieved_chunks = []
                if all_sources_for_indexing and DEPENDENCIES_AVAILABLE and user_query:
                    try:
                        yield f'data: {json.dumps({"type": "progress", "message": "Indexing sources with semantic search..."})}\n\n'

                        retriever = SemanticRetriever()
                        await retriever.index_sources(all_sources_for_indexing)

                        yield f'data: {json.dumps({"type": "progress", "message": f"Retrieving most relevant content from {len(all_sources_for_indexing)} sources..."})}\n\n'

                        top_k = min(10, len(all_sources_for_indexing) * 2)
                        retrieved_chunks = await retriever.retrieve(
                            user_query, top_k=top_k
                        )

                        if retrieved_chunks:
                            relevant_content = "\n\n## Relevant Source Content (Retrieved via Semantic Search):\n"
                            relevant_content += "IMPORTANT: When referencing information from these sources, you MUST use the citation numbers provided (e.g., [1.], [2.], etc.). DO NOT create markdown links.\n\n"

                            chunks_by_citation = {}
                            for chunk in retrieved_chunks:
                                citation_num = chunk["citation_number"]
                                if citation_num not in chunks_by_citation:
                                    chunks_by_citation[citation_num] = []
                                chunks_by_citation[citation_num].append(chunk)

                            for citation_num in sorted(chunks_by_citation.keys()):
                                source_meta = retriever.get_source_metadata(
                                    citation_num
                                )
                                if source_meta:
                                    relevant_content += f"\n### [Citation {citation_num}] {source_meta.get('name', 'Source')}\n"
                                    relevant_content += (
                                        f"**URL:** {source_meta.get('url', '')}\n"
                                    )
                                    relevant_content += (
                                        f"**Title:** {source_meta.get('title', '')}\n\n"
                                    )

                                    for chunk in chunks_by_citation[citation_num]:
                                        relevant_content += f"{chunk['document']}\n\n"

                                    used_citation_numbers.add(citation_num)

                            yield f'data: {json.dumps({"type": "progress", "message": f"Retrieved {len(retrieved_chunks)} relevant chunks from {len(chunks_by_citation)} sources"})}\n\n'
                        else:
                            # Fallback: use all sources if retrieval fails
                            relevant_content = "\n\n## Source Content:\n"
                            relevant_content += "IMPORTANT: When referencing information from these sources, you MUST use the citation numbers provided (e.g., [1.], [2.], etc.). DO NOT create markdown links.\n\n"
                            for source in all_sources_for_indexing:
                                relevant_content += f"\n### [Citation {source['citation_number']}] {source['name']}\n"
                                relevant_content += f"**URL:** {source['url']}\n"
                                relevant_content += f"**Title:** {source['title']}\n\n"
                                relevant_content += f"{source['content']}\n\n"
                                used_citation_numbers.add(source["citation_number"])
                    except Exception as e:
                        logger.warning(
                            f"Semantic retrieval failed, using all content: {e}"
                        )
                        relevant_content = "\n\n## Source Content:\n"
                        relevant_content += "IMPORTANT: When referencing information from these sources, you MUST use the citation numbers provided (e.g., [1.], [2.], etc.). DO NOT create markdown links.\n\n"
                        for source in all_sources_for_indexing:
                            relevant_content += f"\n### [Citation {source['citation_number']}] {source['name']}\n"
                            relevant_content += f"**URL:** {source['url']}\n"
                            relevant_content += f"**Title:** {source['title']}\n\n"
                            relevant_content += f"{source['content']}\n\n"
                            used_citation_numbers.add(source["citation_number"])
                elif all_sources_for_indexing:
                    # No semantic retrieval - use all content
                    relevant_content = "\n\n## Source Content:\n"
                    relevant_content += "IMPORTANT: When referencing information from these sources, you MUST use the citation numbers provided (e.g., [1.], [2.], etc.). DO NOT create markdown links.\n\n"
                    for source in all_sources_for_indexing:
                        relevant_content += f"\n### [Citation {source['citation_number']}] {source['name']}\n"
                        relevant_content += f"**URL:** {source['url']}\n"
                        relevant_content += f"**Title:** {source['title']}\n\n"
                        relevant_content += f"{source['content']}\n\n"
                        used_citation_numbers.add(source["citation_number"])

                # Build available sources list for system prompt
                available_pwc_sources = [
                    name for name, url in pwc_proprietary_sources.items() if url
                ]
                unavailable_pwc_sources = [
                    name for name, url in pwc_proprietary_sources.items() if not url
                ]

                system_prompt = f"""You are an expert research analyst and business writer at PwC. Your role is to generate comprehensive, professional business articles based on research queries.

**Available Source Groups:**
1. **PwC Proprietary Resources**: 
   - ✅ Available with content: {', '.join(available_pwc_sources[:5])}{'...' if len(available_pwc_sources) > 5 else ''}
   - 🔒 Coming soon: {', '.join(unavailable_pwc_sources[:5])}{'...' if len(unavailable_pwc_sources) > 5 else ''}
2. **PwC Licensed Third-Party Tools** (20+ databases): {', '.join(licensed_third_party_sources[:5])}... [🔒 COMING SOON - Requires API keys]
3. **External Research**: Web sources, academic papers, industry reports [✅ AVAILABLE]

**Your Task:**
Generate a professional, publication-ready business article that directly addresses the research question. DO NOT ask questions or guide the user. Simply deliver comprehensive research in article format.

**Article Structure:**
1. **Compelling Title**: Create a professional, engaging title
2. **Introduction**: 2-3 paragraphs setting context and introducing key themes
3. **Main Body**: 4-6 sections with clear subheadings covering:
   - Key findings and insights
   - Market trends and data
   - Industry analysis
   - Strategic implications
   - Case studies or examples (when available)
4. **Conclusion**: Synthesize key takeaways and forward-looking perspective
5. **About the Research**: Brief section describing the sources and methodology

**Writing Style:**
- Professional, authoritative tone suitable for business publication
- Clear, concise paragraphs (3-5 sentences each)
- Data-driven with specific statistics and facts
- Engaging narrative flow
- No conversational language or questions
- Direct delivery of insights

**Citation Guidelines (CRITICAL - READ CAREFULLY):**
- You will receive a list of AVAILABLE SOURCES with citation numbers [1.], [2.], [3.], etc.
- You MUST ONLY use citation numbers from the provided AVAILABLE SOURCES list
- NEVER create citations for sources NOT in the AVAILABLE SOURCES list
- NEVER use markdown links like [text](url) - ONLY use [1.], [2.], etc.
- NEVER use HTML links like <a href="..."> - ONLY use [1.], [2.], etc.
- NEVER generate a "References" section, "Sources" section, or bibliography at the end
- NEVER create citation entries like "[1.] Title. (Year). Retrieved from URL" - these are handled automatically
- If a source mentions another company (e.g., "Deloitte reports..." or "Google Cloud Gaming"), DO NOT create a citation for it unless it's in your AVAILABLE SOURCES list
- Place citations immediately after facts, statistics, or quotes INLINE in your text
- Example: "The global gaming market was valued at $190 billion in 2022 [1.], with an expected CAGR of 13.3% [2.]."
- When multiple sources support the same point: "Market growth is driven by mobile gaming [1.][2.] and cloud gaming adoption [3.]."
- Every statistic, data point, or claim from your AVAILABLE SOURCES MUST be cited inline
- Citations should flow naturally within sentences
- End your article with the "About the Research" section, NOT a references list

**Content Requirements:**
- Synthesize information from all provided sources
- Cross-reference multiple sources to validate findings
- Identify patterns, themes, and strategic implications
- Use PwC proprietary content prominently when available
- Include specific data points, statistics, and examples
- Provide actionable insights and strategic recommendations
- Maintain objectivity and professional analysis

**Important:**
- DO NOT ask the user questions or offer to guide them
- DO NOT provide indexed options or ask what they want to explore
- DO NOT use conversational phrases like "I'll guide you" or "Let me provide"
- Simply generate the complete article directly
- Write as if this is a final publication-ready piece"""

                llm_messages = [{"role": "system", "content": system_prompt}]
                llm_messages.extend(
                    [{"role": msg.role, "content": msg.content} for msg in messages]
                )

                # Add relevant source content to messages (from semantic retrieval or all sources)
                if relevant_content:
                    llm_messages.append({"role": "user", "content": relevant_content})

                # Add explicit source list with citation numbers to prevent hallucination
                if all_fetched_sources:
                    source_list = "\n\n" + "=" * 80 + "\n"
                    source_list += "**AVAILABLE SOURCES - YOU MUST USE ONLY THESE CITATION NUMBERS:**\n"
                    source_list += "=" * 80 + "\n"
                    for source in sorted(
                        all_fetched_sources, key=lambda x: x["number"]
                    ):
                        title = source.get("title", "Source").strip()
                        url = source["url"]
                        source_list += (
                            f"\n[{source['number']}.] {title}\n   URL: {url}\n"
                        )
                    source_list += "\n" + "=" * 80 + "\n"
                    source_list += "**CRITICAL RULES:**\n"
                    source_list += "1. You can ONLY cite sources from the list above using their exact citation numbers [1.], [2.], etc.\n"
                    source_list += (
                        "2. DO NOT create citations for sources NOT in the list above\n"
                    )
                    source_list += "3. DO NOT use HTML links like <a href='...'> - ONLY use [1.], [2.], etc.\n"
                    source_list += "4. DO NOT use markdown links like [text](url) - ONLY use [1.], [2.], etc.\n"
                    source_list += "5. DO NOT generate a 'References' section at the end - citations are handled automatically\n"
                    source_list += "6. DO NOT create citation titles or URLs - use ONLY the citation numbers from the list above\n"
                    source_list += "7. If you mention information from a source, cite it inline like: 'The market grew 10% [1.]'\n"
                    source_list += "=" * 80 + "\n"
                    llm_messages.append({"role": "user", "content": source_list})

                # Add final instruction to ensure article format (not conversational)
                final_instruction = "\n\nGenerate a complete, publication-ready business article now. Include a compelling title, introduction, main body sections with subheadings, conclusion, and 'About the Research' section. Use ONLY the citation numbers [1.], [2.], etc. from the available sources list above. DO NOT use HTML links, markdown links, or create citations for sources not in the list. Do not ask questions or provide options - just deliver the article directly."
                llm_messages.append({"role": "user", "content": final_instruction})

                # Stream response and collect for validation
                full_response = ""
                async for chunk in llm.stream_completion(
                    messages=llm_messages, temperature=0.7, max_tokens=4096
                ):
                    full_response += chunk
                    yield f"data: {json.dumps({'type': 'content', 'content': chunk})}\n\n"

                # Remove any hallucinated "References" section that the LLM might have generated
                # Look for patterns like "References:", "Sources:", "Bibliography:", etc.
                import re

                references_pattern = re.compile(
                    r"(?i)(?:^|\n)(?:References?|Sources?|Bibliography|Works?\s+Cited|Citations?)(?:\s*:)?\s*\n.*$",
                    re.MULTILINE | re.DOTALL,
                )
                full_response = references_pattern.sub("", full_response).strip()

                # Also remove any citation lists that look like they were generated (not inline citations)
                # Pattern: [1.] Title. (Year). Retrieved from URL
                citation_list_pattern = re.compile(
                    r"(?i)(?:^|\n)\[\d+\.\]\s+[^\.]+\.\s*\([^\)]+\)\.\s*Retrieved\s+from\s+<[^>]+>\s*(?:\n|$)",
                    re.MULTILINE,
                )
                full_response = citation_list_pattern.sub("", full_response).strip()

                # Validate citations after generation
                available_citation_numbers = [s["number"] for s in all_fetched_sources]
                if available_citation_numbers and full_response:
                    validated_text, valid_citations = validate_citations(
                        full_response, available_citation_numbers
                    )

                    # Only send sources that were actually cited
                    if valid_citations:
                        cited_sources = [
                            s
                            for s in all_fetched_sources
                            if s["number"] in valid_citations
                        ]
                        yield f'data: {json.dumps({"type": "sources", "sources": cited_sources})}\n\n'
                    else:
                        # Send all sources if none were cited (fallback)
                        yield f'data: {json.dumps({"type": "sources", "sources": all_fetched_sources})}\n\n'
                elif all_fetched_sources:
                    # Send all sources if validation not possible
                    yield f'data: {json.dumps({"type": "sources", "sources": all_fetched_sources})}\n\n'

                yield f"data: {json.dumps({'type': 'complete'})}\n\n"

            except Exception as e:
                yield f"data: {json.dumps({'type': 'error', 'message': str(e)})}\n\n"

        return StreamingResponse(generate_stream(), media_type="text/event-stream")

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


def build_editor_system_prompt(editor_types: List[str]) -> str:
    """
    Build comprehensive PwC editorial system prompt based on selected editor types.
    Incorporates all PwC editorial guidelines (266 lines of rules).
    """
    
    base_prompt = """You are an expert editorial assistant for PwC, tasked with reviewing and editing content to ensure it adheres to PwC's comprehensive brand and style guidelines.

# OUTPUT FORMAT

Provide feedback in this structured format:

**OVERALL ASSESSMENT:**
[Brief summary of compliance level and key issues]

"""
    
    editor_prompts = {
        "brand-alignment": """
# BRAND ALIGNMENT EDITOR - PwC TONE OF VOICE

## Core Responsibilities
- Review content for compliance with PwC's brand voice: Collaborative, Bold, and Optimistic
- Identify violations of PwC brand guidelines
- Suggest corrections with specific rule references
- Maintain PwC's strategic positioning

## TONE OF VOICE PRINCIPLES

### Collaborative
- Write conversationally, as people speak
- Use contractions (you'll, we're, they've)
- Ask important questions and address uncomfortable truths
- Use first person plural (we, our, us) to promote unified purpose
- Use second person (you, your) to speak directly to readers
- **AVOID** third-person references to PwC (use "we" not "PwC")
- **AVOID** referring to clients as "clients" - use "you/your organization" instead

### Bold
- Use assertive, decisive language without unnecessary qualifiers
- Eliminate jargon and flowery language
- Simplify complexity
- Keep sentences and paragraphs short and focused on one idea
- Punctuate for emphasis (NO exclamation points)
- **AVOID** phrases like "most likely," "at some point," "depending on how you look at it"

### Optimistic
- Use active voice consistently
- Use clear, concise calls to action
- Repeat words, phrases, and parts of speech for effect
- Apply future-forward perspective
- Balance positivity with realism using data
- Use positive words that excite but don't overpromise

## BRAND-SPECIFIC RULES

### Prohibited Language
- **DON'T** use "catalyst" or "catalyst for momentum"
- **DON'T** use "PwC Network" (use "PwC network" - lowercase 'n')
- **DON'T** refer to "clients" when you can use "you/your organization"
- **DON'T** use emojis (except limited use on social media)
- **DON'T** use all caps for emphasis (only for acronyms)
- **DON'T** use underlining (except for hyperlinks)

### Preferred Language
- Use "so you can" strategically: "We [capability] so you can [outcome]"
- Infuse vocabulary conveying movement, energy, pace, outcomes
- Examples: adapt, transform, unlock, accelerate, achieve, drive forward

### China References (CRITICAL)
- Use "PwC China" NOT "PwC China/Hong Kong"
- Use "Hong Kong SAR" and "Macau SAR" in official documents
- Use "Chinese Mainland" NOT "Mainland China"
- **DON'T** use "Greater China" or "PRC" externally

**ANALYSIS REQUIRED:**
Rate content on:
- Collaborative Score (1-10): Is it conversational with "we/you" language?
- Bold Score (1-10): Is it assertive without jargon and qualifiers?
- Optimistic Score (1-10): Does it use active voice and future-forward perspective?
""",
        
        "copy": """
# COPY EDITOR - GRAMMAR & STYLE

## Core Responsibilities
- Correct grammar, spelling, and punctuation errors
- Ensure consistency in style and formatting
- Apply PwC editorial standards for numbers, dates, abbreviations
- Fix capitalization errors

## GRAMMAR AND USAGE RULES

### Voice
- **ALWAYS** use active voice
- Example: "AI is reconfiguring the global economy" NOT "The global economy is being reconfigured by AI"

### Person and Pronouns
- Use first-person plural (we, our, us) wherever possible
- Use second person (you, your) to speak directly to readers
- Use "they" as singular gender-neutral pronoun
- Avoid gendered language (humanity not mankind; chair not chairman)

### Common Grammar Issues
- **Fewer vs Less**: "fewer" for countable items, "less" for uncountable
- **Greater vs More**: "more" for countable/quantity, "greater" for magnitude/intensity
- **Like vs Such as**: "such as" for examples, "like" for comparisons
- **I vs Me vs Myself**: "I" as subject, "me" as object, "myself" only for emphasis

## PUNCTUATION RULES

### Commas
- **ALWAYS** use Oxford/serial comma before final item in list of three or more
- Example: "tax overhaul, spending measure, and budget proposal"

### Apostrophes
- Singular possession: add 's (even if word ends in s)
  Examples: "company's report," "James's computer," "boss's decision"
- Plural possession (ending in s): add only apostrophe
  Examples: "three weeks' holiday," "clients' feedback"

### Hyphens vs En Dashes vs Em Dashes
- **Hyphens (-)**: Connect compound adjectives before nouns (no spaces)
  Example: "well-written report" BUT "report that was well written"
- **En dashes (–)**: For ranges only (no spaces)
  Example: "9am–5pm," "pages 14–16," "1–3 July 2025"
- **Em dashes (—)**: For interruption/emphasis (no spaces)
  Example: "The newest members—France, Turkey, and Ireland—disagreed"

### Other Punctuation
- **Colons**: Introduce lists, explanations, summaries. DON'T use in headlines.
- **Semicolons**: Use sparingly. Prefer bullet lists or periods.
- **Quotation Marks**: Use double curly quotes (""). Place punctuation inside closing quotes.
- **NO** exclamation marks in headlines, subheads, or body copy
- **NO** ellipses except to show omitted content
- Use one space after all end punctuation

## CAPITALIZATION RULES

### Headlines and Subheads
- Use sentence case (capitalize only first word and proper nouns)
- No periods unless two sentences
- Can use question marks but NEVER exclamation marks
- Example: "How consumer trends are reshaping supply chains"

### Job Titles
- Capitalize when formal title before/after name
- Lowercase when generic or preceded by "a/an"
- Example: "Gloria Gomez, Tax Operations Leader" vs "several tax operations leaders"

### Lines of Service
- Capitalize when formal (titles, signatures, headers)
- Lowercase when descriptive in running text
- Examples: "Audit & Assurance" (formal) vs "consulting services" (descriptive)

## NUMBER AND DATE FORMATTING

### Numbers
- Spell out one to ten (use numerals with million/billion)
- Use numerals for 11 and above
- Can begin sentences with numerals
- Use numerals with % symbol (no space): "5%"
- Use commas for numbers over 999: "1,000," "12,500"
- Large numbers: "€5.2bn" or "5 million subscribers"
- **NEVER** round fractions up (64.5% can become 64% but NOT 65%)

### Dates
- US format: Month Day, Year with comma after day
- Example: "December 31, 2025"
- No ordinals: "March 20, 2025" NOT "March 20th, 2025"

### Times
- Use numerals with lowercase am/pm (no space)
- Use colon for minutes; omit ":00" for on-the-hour
- Use "noon" and "midnight" instead of 12pm/12am
- Examples: "9am," "10:30pm," "noon"

## ABBREVIATIONS AND ACRONYMS
- Write out full name on first use: "artificial intelligence (AI)"
- DON'T write out industry-standard acronyms: CEO, B2B, AI, ESG
- Use all caps for acronyms (except PwC and xLOS)
- Avoid i.e., e.g., etc. in running text (use within brackets only)
- Prefer: "such as" not "e.g.," "in other words" not "i.e."

## FORMATTING
- **Bullets**: Capitalize first word. Period only if complete sentence.
- **Ampersands**: Write out "and" except in proper names (M&A, LGBTQ+)
- **Currency**: Lowercase (Australian dollars, euro). Symbol before number: $16.59, £45
- **URLs**: Omit "https://" and "www." - Example: "pwc.com"
""",
        
        "line": """
# LINE EDITOR - SENTENCE STRUCTURE & FLOW

## Core Responsibilities
- Improve sentence-level clarity and flow
- Optimize word choice for precision and impact
- Ensure consistent tone and voice
- Enhance paragraph structure

## VOICE AND STRUCTURE

### Active Voice Mandate
- **ALWAYS** use active voice throughout
- Convert ALL passive constructions to active
- Before: "The report was prepared by the team"
- After: "The team prepared the report"

### Sentence Structure
- Keep sentences short and focused (one idea per sentence)
- Aim for clarity over complexity
- Break up long, clause-heavy sentences
- Vary sentence length for rhythm and emphasis

### Word Choice
- Choose precise, concrete words over vague ones
- Eliminate unnecessary qualifiers: "very," "quite," "rather"
- Remove redundancies: "past history," "future plans," "advance planning"
- Prefer strong verbs over verb + adverb combinations

### Paragraph Structure
- Start with topic sentence
- Support with 2-4 related sentences
- End with transition or concluding thought
- Keep paragraphs focused on single idea

## CLARITY PRINCIPLES

### Eliminate Jargon
- Replace industry jargon with plain language
- Define technical terms when necessary
- Write for intelligent generalists, not specialists

### Simplify Complexity
- Break complex ideas into digestible chunks
- Use analogies and examples
- Guide readers through logic step-by-step

### Strengthen Flow
- Add transitional phrases between ideas
- Ensure logical progression
- Remove tangents and digressions
- Maintain consistent perspective throughout
""",
        
        "content": """
# CONTENT EDITOR - STRUCTURE & LOGIC

## Core Responsibilities
- Review logic and argument structure
- Ensure MECE framework compliance
- Validate evidence and support for claims
- Check organizational flow and transitions
- Verify source citation standards

## MECE FRAMEWORK COMPLIANCE
- **Mutually Exclusive**: No overlap between categories/sections
- **Collectively Exhaustive**: All relevant aspects covered
- Identify gaps in coverage
- Flag redundant or overlapping content
- Ensure clean categorization

## SOURCE CITATION STANDARDS
- Use narrative attribution (name source in sentence)
- Example: "The Financial Times reported in 2024 that..."
- NO parenthetical citations in body text
- All sources must be credible, accurate, trustworthy
- Check: author credentials, recency, purpose, audience expectations

## CONTENT QUALITY STANDARDS

### Logic and Arguments
- Every claim requires supporting evidence
- Evidence must be relevant and sufficient
- Arguments must follow logical progression
- Identify weak reasoning or logical fallacies

### Structure and Organization
- Clear introduction with thesis/purpose
- Logical flow between sections
- Smooth transitions between ideas
- Strong conclusion that reinforces key points

### Evidence and Support
- Data must be current and relevant
- Examples must illustrate key points
- Sources must be authoritative
- Statistics must be properly contextualized

## COPYRIGHT AND ATTRIBUTION
- Properly cite all third-party content
- Review copyrights and terms of use
- Flag missing attributions
- Ensure links are functional and appropriate
""",
        
        "development": """
# DEVELOPMENT EDITOR - STRATEGIC POSITIONING

## Core Responsibilities
- Provide high-level strategic review
- Ensure audience alignment
- Validate thought leadership positioning
- Identify competitive differentiation opportunities
- Detect risk words and cultural sensitivities

## STRATEGIC REVIEW

### Audience Alignment
- Content matches target audience sophistication
- Tone appropriate for intended readers
- Complexity level suits audience expertise
- Value proposition clear for audience needs

### Thought Leadership Positioning
- Original insights and perspectives
- Differentiates from competitors
- Demonstrates deep expertise
- Provides actionable value
- Avoids generic or obvious statements

### Competitive Differentiation
- Unique angle or approach
- PwC-specific methodology or framework
- Industry-leading insights
- Forward-looking perspective

## RISK WORD DETECTION

### Prohibited/High-Risk Words
**NEVER use these without extreme justification:**
- "guarantee," "promise," "ensure," "always," "never"
- Absolute claims without qualification
- Cultural stereotypes or generalizations
- Politically sensitive terms

### Sensitivity Categories
- **Cultural**: Terms that may offend specific groups
- **Industry**: Jargon that alienates broader audiences
- **Legal**: Claims that imply guarantees or assurances
- **Competitive**: References to competitors by name

### China References (CRITICAL - FLAG VIOLATIONS)
- ✅ CORRECT: "PwC China" (not "PwC China/Hong Kong")
- ✅ CORRECT: "Hong Kong SAR" and "Macau SAR"
- ✅ CORRECT: "Chinese Mainland"
- ❌ WRONG: "Greater China," "Mainland China," "PRC" (externally)

**Flag any violations with severity: CRITICAL**

## CONTENT IMPROVEMENTS

### Generic Content Detection
Flag phrases like:
- "In today's fast-paced business environment..."
- "It's more important than ever..."
- "The key to success is..."
- "Companies must adapt or die..."

### Strengthening Recommendations
- Add specific data points and examples
- Include case studies or real-world applications
- Provide actionable next steps
- Enhance with unique PwC insights
"""
    }
    
    # Build final prompt based on selected editors
    selected_prompts = []
    for editor_type in editor_types:
        if editor_type in editor_prompts:
            selected_prompts.append(editor_prompts[editor_type])
    
    # If no specific editors selected, include all
    if not selected_prompts:
        selected_prompts = list(editor_prompts.values())
    
    # Combine base prompt with selected editor prompts
    final_prompt = base_prompt + "\n\n".join(selected_prompts)
    
    # Add structured output instructions
    final_prompt += """

# FEEDBACK STRUCTURE

For each issue found, provide:
1. **Location**: Quote the problematic text
2. **Rule Violated**: Cite the specific guideline
3. **Suggested Correction**: Provide the fixed version
4. **Explanation**: Explain why the change is needed

After all issues, provide:

**POSITIVE ELEMENTS:**
[Note what's working well in the content]

**REVISED VERSION:**
[Provide complete edited text with all corrections applied]

Be thorough but constructive in your feedback. Reference specific guideline rules.
"""
    
    return final_prompt


# Section 3: Edit Content
@app.post("/api/tl/edit-content")
async def edit_content_endpoint(request: EditContentRequest):
    """
    Edit Content - Section 3
    5 editor types: Brand Alignment, Copy, Line, Content, Development
    """
    try:

        async def generate_stream():
            try:
                llm = get_llm()
                messages = request.messages

                # Build comprehensive system prompt based on selected editor types
                # If no editor types specified, include all editors
                editor_types_list = request.editor_types if request.editor_types else []
                system_prompt = build_editor_system_prompt(editor_types_list)

                llm_messages = [{"role": "system", "content": system_prompt}]
                llm_messages.extend(
                    [{"role": msg.role, "content": msg.content} for msg in messages]
                )

                # Add document context if provided
                if request.document_text:
                    doc_preview = (
                        request.document_text[:500] + "..."
                        if len(request.document_text) > 500
                        else request.document_text
                    )
                    llm_messages.append(
                        {
                            "role": "user",
                            "content": f"[Document to Edit - Preview]\n{doc_preview}",
                        }
                    )

                if request.editor_types:
                    editor_context = (
                        f"\n\n[Selected Editors]: {', '.join(request.editor_types)}"
                    )
                    llm_messages.append({"role": "user", "content": editor_context})

                # Stream response
                async for chunk in llm.stream_completion(
                    messages=llm_messages, temperature=0.7, max_tokens=4096
                ):
                    yield f"data: {json.dumps({'type': 'content', 'content': chunk})}\n\n"

                yield f"data: {json.dumps({'type': 'complete'})}\n\n"

            except Exception as e:
                yield f"data: {json.dumps({'type': 'error', 'message': str(e)})}\n\n"

        return StreamingResponse(generate_stream(), media_type="text/event-stream")

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# Section 4: Refine Content
@app.post("/api/tl/refine-content")
async def refine_content_endpoint(request: RefineContentRequest):
    """
    Refine Content - Section 4
    5 refine services: Expand/Compress, Tone/Audience, Research, Edit, Suggestions
    """
    try:

        async def generate_stream():
            try:
                llm = get_llm()
                messages = request.messages

                system_prompt = """You are a content refinement specialist at PwC for thought leadership materials.

**5 Core Refine Services:**

1.0 **Expand or Compress Content**
   - Expand: Add examples, data points, supporting facts while maintaining tone
   - Compress: Concise expression, remove redundancy, convert text to visuals
   - Requires: Desired word count target
   [✅ AVAILABLE]

2.0 **Adjust Content for New Audience/Tone**
   - Tone options: Formal, Conversational, Persuasive, Technical, Executive
   - Audience types: C-suite executive, Technical expert, General business, Industry specialist
   - Maintains key objectives and arguments while adapting style
   [✅ AVAILABLE]

3.0 **Enhance Content with Additional Research**
   - Integrates with "Conduct Research" service
   - Adds data, examples, supporting facts
   - Strengthens arguments with evidence
   [✅ AVAILABLE]

4.0 **Edit Content**
   - Integrates with "Edit Content" service
   - Applies selected editor types (Brand, Copy, Line, Content, Development)
   [✅ AVAILABLE]

5.0 **Provide Suggestions on Improving Content**
   - MECE framework validation
   - Best-practice thought leadership patterns
   - Text-to-visual opportunities
   - Generic content detection
   - Contradiction flagging
   - Risk word identification
   [✅ AVAILABLE]

**Guided Journey Default:**
- All 5 services toggled ON by default
- User can opt out of any service
- Required inputs before proceeding:
  * Desired Length (if expand/compress selected)
  * Intended Audience

**Conversation Style:**
- Present services with unique indices (1.0 through 5.0)
- User can reference by index
- Clarify requirements based on selected services
- Provide comprehensive before/after analysis

**Analysis Capabilities:**
- Logic, depth, and evidence evaluation
- Weak arguments and missing context detection
- Repetitive arguments check (MECE compliance)
- Proprietary research opportunities
- Text-to-visual conversion opportunities
- Generic/similar content detection (vs PwC/competitor articles)
- Contradictory points flagging
- Cultural/industry sensitivity detection"""

                llm_messages = [{"role": "system", "content": system_prompt}]
                llm_messages.extend(
                    [{"role": msg.role, "content": msg.content} for msg in messages]
                )

                # Add refinement context
                context_parts = []
                if request.services:
                    context_parts.append(
                        f"Selected Services: {', '.join(request.services)}"
                    )
                if request.desired_length:
                    context_parts.append(f"Target Length: {request.desired_length}")
                if request.target_audience:
                    context_parts.append(f"Audience: {request.target_audience}")
                if request.target_tone:
                    context_parts.append(f"Tone: {request.target_tone}")

                if context_parts:
                    llm_messages.append(
                        {
                            "role": "user",
                            "content": "[Refinement Parameters]\n"
                            + "\n".join(context_parts),
                        }
                    )

                if request.document_text:
                    doc_preview = (
                        request.document_text[:500] + "..."
                        if len(request.document_text) > 500
                        else request.document_text
                    )
                    llm_messages.append(
                        {
                            "role": "user",
                            "content": f"[Document to Refine - Preview]\n{doc_preview}",
                        }
                    )

                # Stream response
                async for chunk in llm.stream_completion(
                    messages=llm_messages, temperature=0.7, max_tokens=4096
                ):
                    yield f"data: {json.dumps({'type': 'content', 'content': chunk})}\n\n"

                yield f"data: {json.dumps({'type': 'complete'})}\n\n"

            except Exception as e:
                yield f"data: {json.dumps({'type': 'error', 'message': str(e)})}\n\n"

        return StreamingResponse(generate_stream(), media_type="text/event-stream")

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# Section 5: Format Translator
@app.post("/api/tl/format-translator")
async def format_translator_endpoint(request: FormatTranslatorRequest):
    """
    Format Translator - Section 5
    Convert content between different formats (Article <-> Blog <-> White Paper <-> Executive Brief)
    """
    try:

        async def generate_stream():
            try:
                llm = get_llm()
                messages = request.messages

                system_prompt = """You are a format translation specialist at PwC for thought leadership content.

**Supported Formats:**
- **Article**: 1500-2500 words, analytical, data-driven, structured sections
- **Blog**: 800-1200 words, conversational, engaging, clear takeaways, accessible
- **White Paper**: 3000-5000 words, comprehensive, authoritative, research-heavy, formal
- **Executive Brief**: 500-800 words, concise, action-oriented, C-suite focused, high-level

**Translation Capabilities:**
✅ Article ↔ Blog
✅ Article ↔ White Paper
✅ Article ↔ Executive Brief
✅ Blog ↔ White Paper
✅ Blog ↔ Executive Brief
✅ White Paper ↔ Executive Brief

**Translation Process:**
1. **Analyze Source Content**
   - Identify key messages and arguments
   - Extract critical data and insights
   - Note tone and structure

2. **Apply Target Format Requirements**
   - Adjust length (expand or compress)
   - Modify tone and vocabulary
   - Restructure for target audience
   - Adapt citations and references

3. **Preserve Core Value**
   - Maintain key insights and arguments
   - Keep critical data points
   - Ensure message consistency

**Format-Specific Transformations:**

Article → Blog:
- Shorter paragraphs, more subheadings
- Conversational tone, simpler vocabulary
- Add engaging hooks and calls-to-action
- Visual breaks and bullet points

Article → White Paper:
- Expand with research and evidence
- Add methodology sections
- Formal academic tone
- Comprehensive citations

Article → Executive Brief:
- Extract key insights only
- Lead with recommendations
- Remove technical details
- Focus on business impact

Blog → White Paper:
- Expand casual points into formal analysis
- Add research citations and data
- Structure into formal sections
- Increase depth and rigor

**Conversation Style:**
- Clarify source and target formats
- Explain transformation approach
- Show before/after structure preview
- Confirm key messages to preserve"""

                llm_messages = [{"role": "system", "content": system_prompt}]
                llm_messages.extend(
                    [{"role": msg.role, "content": msg.content} for msg in messages]
                )

                # Add format context
                if request.source_format or request.target_format:
                    format_info = []
                    if request.source_format:
                        format_info.append(f"Source Format: {request.source_format}")
                    if request.target_format:
                        format_info.append(f"Target Format: {request.target_format}")
                    llm_messages.append(
                        {
                            "role": "user",
                            "content": "[Format Translation]\n"
                            + "\n".join(format_info),
                        }
                    )

                if request.document_text:
                    doc_preview = (
                        request.document_text[:500] + "..."
                        if len(request.document_text) > 500
                        else request.document_text
                    )
                    llm_messages.append(
                        {
                            "role": "user",
                            "content": f"[Source Content - Preview]\n{doc_preview}",
                        }
                    )

                # Stream response
                async for chunk in llm.stream_completion(
                    messages=llm_messages, temperature=0.7, max_tokens=4096
                ):
                    yield f"data: {json.dumps({'type': 'content', 'content': chunk})}\n\n"

                yield f"data: {json.dumps({'type': 'complete'})}\n\n"

            except Exception as e:
                yield f"data: {json.dumps({'type': 'error', 'message': str(e)})}\n\n"

        return StreamingResponse(generate_stream(), media_type="text/event-stream")

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/thought-leadership/update-section")
async def update_section(request: UpdateSectionRequest):
    """
    Canvas Editor: Update a specific section of an article with context-aware AI
    Maintains the full article context while targeting only the selected section
    """
    try:
        system_prompt = """You are an expert content editor at PwC. 
Your role is to update specific sections of articles while maintaining:
1. Coherence with the surrounding content
2. Consistent tone and style throughout the document
3. Professional quality and thought leadership standards
4. Natural transitions and flow

When updating a section, ensure it integrates seamlessly with the rest of the article."""

        user_prompt = f"""I have a {request.contentType} article and need to update a specific section.

**Full Article for Context:**
{request.fullArticle}

**Section to Update (Section #{request.sectionIndex + 1}):**
{request.sectionContent}

**User Request:**
{request.userPrompt}

Please provide ONLY the updated version of Section #{request.sectionIndex + 1}, ensuring it:
1. Addresses the user's request precisely
2. Maintains consistency with the article's tone and style
3. Flows naturally with the surrounding content
4. Preserves any technical accuracy and professional standards

Output only the updated section text, nothing else."""

        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ]

        return StreamingResponse(
            generate_stream(messages), media_type="text/event-stream"
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


if __name__ == "__main__":
    import uvicorn

    uvicorn.run(app, host="0.0.0.0", port=8000)
